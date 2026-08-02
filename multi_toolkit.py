#!/usr/bin/env python3
"""
Multi Toolkit — a local, cross-platform workbench for PDFs, media & QR codes.

The browser is the interface, so there is no GUI toolkit to install. A tiny
standard-library web server does the work; everything runs on your machine and
nothing is uploaded anywhere.

PDF
---
  • Merge      — combine many PDFs (drag rows to reorder), optional bookmarks
  • Split      — VISUAL page preview: click between pages to place cut points,
                 or click pages to extract them; plus every-page / every-N /
                 text ranges / reorder
  • Compress   — lossless cleanup, or strong image recompression/downscaling
  • PDF → MD   — Markdown for pasting into an LLM: text stays text (a fraction
                 of the tokens a PDF costs), figures are CROPPED to their own
                 bounding box and anchored at the point in the reading order
                 where they occur, each above its own caption. Column-aware for
                 one, two and three-column layouts. Tables become real Markdown
                 tables, display equations are fenced in $$, running heads are
                 dropped. A page with no text layer is exported whole and
                 flagged rather than silently mangled. Every .md opens with a
                 short note explaining its own layout, so it can be handed to a
                 model with no covering message. Also available headless:
                     python multi_toolkit.py md <pdf-or-folder> [outdir]
  • Convert    — Office → PDF (Word/PowerPoint/Excel)   [needs LibreOffice]
                 PDF → Word (.docx)
                 PDF → PowerPoint (.pptx, page-image slides)
                 PDF → Images (PNG/JPG, zipped)
                 Images → PDF (combined, in list order)

Photo & Video
-------------
  • YouTube    — paste a URL, pick resolution + format, download with live
                 progress (%, speed, ETA). Uses yt-dlp.
                 Formats: "QuickTime mp4" downloads max-res then re-encodes
                 to HEVC/H.264 (hardware-accelerated where available) so 4K
                 plays natively on macOS; "Original codec" keeps YouTube's
                 VP9/AV1 (max quality — use VLC/IINA); "Native H.264" needs
                 no conversion but YouTube caps it at ~1080p.
  • Reels      — the same engine pointed at Instagram / TikTok / Shorts, plus
                 9:16 reframing: blurred backdrop (nothing cropped), crop to
                 fill, or black bars. Canvas 1080x1920, 720x1280 or 1080x1350.
                 Can borrow cookies from an installed browser for posts that
                 need a signed-in session.
  • Image      — load from a URL or a local file, drag a crop box (free or
                 locked to 1:1 / 4:5 / 9:16 / 16:9 / 3:2), rotate, flip,
                 grayscale, resize, then export PNG / JPG / WebP.
  • QR Code    — links, plain text, Wi-Fi, email, SMS, phone, contact card or
                 coordinates. Square / dot / rounded modules, any colours,
                 transparent background, and a CENTRE LOGO — drop in a PNG and
                 optionally flatten it to a silhouette so it reads as part of
                 the code. Live preview; error correction is forced to H when
                 a logo is present, and the app warns when a choice is likely
                 to hurt scannability.
                 ffmpeg is required for anything above ~720p.
                 Only download videos you own or have permission to save.

Quality-of-life
---------------
  • Files are uploaded once and cached server-side under a token, so repeated
    operations on big PDFs don't re-send the bytes.
  • Every PDF row shows a first-page thumbnail and a page count.
  • Drag rows to reorder; image files show their own previews.

Dependencies
------------
Core (always): pypdf — auto-installed on first run.
Installed automatically the first time a feature needs them (pip wheels, no
system tools required): pypdfium2, python-pptx, pdf2docx, pikepdf, Pillow,
yt-dlp, qrcode, PyMuPDF.

LibreOffice (OPTIONAL, only for Office→PDF): the one piece that can't be a pip
wheel. Free, cross-platform. If it's missing the app still runs and tells you
how to get it:  https://www.libreoffice.org/download

ffmpeg (OPTIONAL, needed for max-resolution YouTube and for all 9:16
reframing): YouTube serves >720p video and audio as separate streams; ffmpeg
merges them. macOS: brew install ffmpeg
Windows: winget install Gyan.FFmpeg   Linux: apt install ffmpeg

Run
---
    python multi_toolkit.py

Opens in your default browser automatically. Ctrl-C in the terminal to stop.
"""

import base64
import importlib
import io
import json
import os
import re
import shutil
import socket
import subprocess
import time
import sys
import tempfile
import threading
import traceback
import uuid
import zipfile
import webbrowser
from collections import OrderedDict
from http.server import BaseHTTPRequestHandler, ThreadingHTTPServer
from urllib.parse import urlparse


# --------------------------------------------------------------------------- #
#  Dependency bootstrap
# --------------------------------------------------------------------------- #
_INSTALL_LOCK = threading.Lock()


def ensure(pkg, import_name=None):
    """Import a package, installing it via pip on demand. Returns module or None."""
    name = import_name or pkg
    try:
        return importlib.import_module(name)
    except ImportError:
        pass
    with _INSTALL_LOCK:
        try:  # re-check inside the lock (another thread may have installed it)
            return importlib.import_module(name)
        except ImportError:
            pass
        print(f"[setup] installing '{pkg}' ...")
        for extra in ([], ["--user"], ["--break-system-packages"]):
            try:
                subprocess.check_call(
                    [sys.executable, "-m", "pip", "install", "--quiet", *extra, pkg]
                )
                importlib.invalidate_caches()
                return importlib.import_module(name)
            except Exception:  # noqa: BLE001
                continue
        print(f"[setup] could not install '{pkg}'.")
        return None


class FeatureError(RuntimeError):
    """User-facing error (sent back as a readable message, not a traceback)."""


class CacheMiss(RuntimeError):
    """A file token is no longer cached; the client should resend the bytes."""


if ensure("pypdf") is None:
    sys.exit("ERROR: 'pypdf' is required and could not be installed automatically.\n"
             "Try:  pip install pypdf")
from pypdf import PdfReader, PdfWriter  # noqa: E402


def need(pkg, import_name=None, hint=None):
    mod = ensure(pkg, import_name)
    if mod is None:
        raise FeatureError(hint or f"Could not install '{pkg}'. Try: pip install {pkg}")
    return mod


def office_binary():
    return shutil.which("soffice") or shutil.which("libreoffice")


def gs_binary():
    return shutil.which("gs") or shutil.which("gswin64c") or shutil.which("gswin32c")


def ffmpeg_binary():
    return shutil.which("ffmpeg")


def safe_stem(name):
    stem = os.path.splitext(os.path.basename(name or "file"))[0]
    stem = re.sub(r"[^\w.\- ]+", "_", stem).strip() or "file"
    return stem


# --------------------------------------------------------------------------- #
#  Server-side file cache (upload once, operate many times)
# --------------------------------------------------------------------------- #
_CACHE_LOCK = threading.Lock()
_FILE_CACHE: "OrderedDict[str, tuple[str, bytes]]" = OrderedDict()
_CACHE_MAX_FILES = 64
_CACHE_MAX_BYTES = 600 * 1024 * 1024  # 600 MB


def cache_put(name, data):
    token = uuid.uuid4().hex
    with _CACHE_LOCK:
        _FILE_CACHE[token] = (name, data)
        # Evict oldest entries beyond the caps.
        while (len(_FILE_CACHE) > _CACHE_MAX_FILES
               or sum(len(d) for _, d in _FILE_CACHE.values()) > _CACHE_MAX_BYTES):
            if len(_FILE_CACHE) <= 1:
                break
            _FILE_CACHE.popitem(last=False)
    return token


def cache_get(token):
    with _CACHE_LOCK:
        item = _FILE_CACHE.get(token)
        if item is not None:
            _FILE_CACHE.move_to_end(token)  # LRU touch
        return item


def resolve_files(items):
    """Turn client file refs ({token} or {name,data}) into {name, bytes} dicts.

    If a file arrives with both token and data, the data wins and refreshes
    the cache. A token that's no longer cached raises CacheMiss (HTTP 409) so
    the client can transparently resend the bytes.
    """
    out = []
    for it in items or []:
        name = it.get("name", "file")
        if it.get("data"):
            data = base64.b64decode(it["data"])
            token = cache_put(name, data)
        elif it.get("token"):
            hit = cache_get(it["token"])
            if hit is None:
                raise CacheMiss(name)
            name, data = hit
            token = it["token"]
        else:
            raise FeatureError(f"File entry for “{name}” has neither data nor token.")
        out.append({"name": name, "bytes": data, "token": token})
    return out


# --------------------------------------------------------------------------- #
#  Inspect: page counts + thumbnails for the visual UI
# --------------------------------------------------------------------------- #
_THUMB_CAP = 600          # pages beyond this render as number-only cards
_THUMB_WIDTH = 150        # px width of split-grid thumbnails
_COVER_WIDTH = 96         # px width of file-row cover thumbnails


def _render_thumb(page, width):
    w_pt, _h_pt = page.get_size()
    scale = max(0.05, width / max(1.0, w_pt))
    pil = page.render(scale=scale).to_pil()
    if pil.mode != "RGB":
        pil = pil.convert("RGB")
    buf = io.BytesIO()
    pil.save(buf, "JPEG", quality=72)
    return base64.b64encode(buf.getvalue()).decode()


def op_inspect(file, thumbs="cover"):
    """Return page count and thumbnail(s) for a PDF.

    thumbs: 'none' | 'cover' (first page only) | 'all' (every page, capped)
    """
    pdfium = need("pypdfium2")
    try:
        doc = pdfium.PdfDocument(file["bytes"])
    except Exception:  # noqa: BLE001
        raise FeatureError(f"“{file['name']}” doesn't look like a readable PDF.")
    n = len(doc)
    result = {"pages": n, "token": file["token"], "cover": None, "thumbs": None}
    try:
        if thumbs == "cover" and n:
            result["cover"] = _render_thumb(doc[0], _COVER_WIDTH)
        elif thumbs == "all" and n:
            imgs = []
            for i in range(min(n, _THUMB_CAP)):
                imgs.append(_render_thumb(doc[i], _THUMB_WIDTH))
            result["thumbs"] = imgs
            result["cover"] = imgs[0]
    finally:
        doc.close()
    return result


# --------------------------------------------------------------------------- #
#  Core operations  (each returns a list of (filename, bytes), info string)
# --------------------------------------------------------------------------- #
def op_merge(files, bookmarks=True):
    writer = PdfWriter()
    cursor = 0
    for f in files:
        reader = PdfReader(io.BytesIO(f["bytes"]))
        if reader.is_encrypted:
            try:
                reader.decrypt("")
            except Exception:  # noqa: BLE001
                pass
        start = cursor
        for page in reader.pages:
            writer.add_page(page)
            cursor += 1
        if bookmarks and cursor > start:
            try:
                writer.add_outline_item(safe_stem(f["name"]), start)
            except Exception:  # noqa: BLE001
                pass
    buf = io.BytesIO()
    writer.write(buf)
    return [("merged.pdf", buf.getvalue())], \
        f"Merged {len(files)} files ({cursor} pages) into merged.pdf"


def _parse_ranges(spec, total):
    out = []
    for chunk in str(spec).split(","):
        chunk = chunk.strip()
        if not chunk:
            continue
        if "-" in chunk:
            a, b = chunk.split("-", 1)
            a = int(a) if a.strip() else 1
            b = int(b) if b.strip() else total
        else:
            a = b = int(chunk)
        a = max(1, a)
        b = min(total, b)
        if a <= b:
            out.append((a, b))
    if not out:
        raise FeatureError("No valid page ranges were parsed.")
    return out


def op_split(file, mode="each", ranges="", every_n=1, pages=None):
    reader = PdfReader(io.BytesIO(file["bytes"]))
    if reader.is_encrypted:
        try:
            reader.decrypt("")
        except Exception:  # noqa: BLE001
            pass
    total = len(reader.pages)
    stem = safe_stem(file["name"])
    results = []

    def write_pages(idxs, label):
        w = PdfWriter()
        for i in idxs:
            w.add_page(reader.pages[i])
        b = io.BytesIO()
        w.write(b)
        results.append((f"{stem}_{label}.pdf", b.getvalue()))

    if mode == "each":
        for i in range(total):
            write_pages([i], f"p{i + 1:03d}")
        info = f"Split into {total} single-page PDFs."
    elif mode == "ranges":
        for (a, b) in _parse_ranges(ranges, total):
            write_pages(range(a - 1, b), f"{a}-{b}" if a != b else f"p{a}")
        info = f"Split into {len(results)} PDFs."
    elif mode == "every":
        n = max(1, int(every_n))
        for start in range(0, total, n):
            end = min(start + n, total)
            write_pages(range(start, end), f"{start + 1}-{end}")
        info = f"Split into {len(results)} chunks of up to {n} pages."
    elif mode == "extract":
        sel = sorted({int(p) for p in (pages or []) if 1 <= int(p) <= total})
        if not sel:
            raise FeatureError("Select at least one page to extract.")
        write_pages([p - 1 for p in sel], "extract")
        info = f"Extracted {len(sel)} page(s) into one PDF."
    elif mode == "reorder":
        seq = [int(p) for p in (pages or []) if 1 <= int(p) <= total]
        if not seq:
            raise FeatureError("No pages left to write — restore at least one page.")
        write_pages([p - 1 for p in seq], "reordered")
        dropped = total - len(set(seq))
        info = f"Rebuilt PDF with {len(seq)} page(s) in new order" + \
               (f" ({dropped} removed)." if dropped > 0 else ".")
    else:
        raise FeatureError(f"Unknown split mode: {mode}")
    return results, info


COMPRESS_PRESETS = {
    "high":     {"mode": "strong",   "quality": 35, "max_px": 1100},
    "balanced": {"mode": "strong",   "quality": 60, "max_px": 1600},
    "quality":  {"mode": "strong",   "quality": 85, "max_px": 2400},
    "lossless": {"mode": "lossless", "quality": 0,  "max_px": 0},
}


def op_compress_preview(files):
    """Compress every file with every preset; cache results, return sizes."""
    pikepdf = need("pikepdf")
    out = {"orig": sum(len(f["bytes"]) for f in files), "presets": {}}
    for key, p in COMPRESS_PRESETS.items():
        entries, total = [], 0
        for f in files:
            data = _compress_one(pikepdf, f["bytes"], p["mode"], p["quality"], p["max_px"])
            if len(data) >= len(f["bytes"]):
                data = f["bytes"]  # already optimal — keep original
            name = f"{safe_stem(f['name'])}_compressed.pdf"
            entries.append({"name": name, "size": len(data),
                            "orig": len(f["bytes"]), "token": cache_put(name, data)})
            total += len(data)
        out["presets"][key] = {"total": total, "files": entries}
    return out


def op_compress(files, mode="strong", quality=60, max_px=1600):
    pikepdf = need("pikepdf")
    results = []
    notes = []
    for f in files:
        data = f["bytes"]
        orig = len(data)
        out = _compress_one(pikepdf, data, mode, quality, max_px)
        # If "compression" made it bigger (already optimized), keep the smaller.
        if len(out) >= orig:
            out = data
            saved = "already optimal — kept original"
        else:
            saved = f"{orig // 1024} KB -> {len(out) // 1024} KB " \
                    f"({100 * len(out) / max(1, orig):.0f}%)"
        results.append((f"{safe_stem(f['name'])}_compressed.pdf", out))
        notes.append(f"{safe_stem(f['name'])}: {saved}")
    return results, "\n".join(notes)


def _compress_one(pikepdf, data, mode, quality, max_px):
    from pikepdf import Pdf, Name
    # Strong mode: prefer Ghostscript if present (best general result).
    if mode == "strong" and gs_binary():
        try:
            return _gs_compress(data, quality)
        except Exception:  # noqa: BLE001
            pass  # fall through to pikepdf path

    pdf = Pdf.open(io.BytesIO(data))
    if mode == "strong":
        PdfImage = pikepdf.PdfImage
        need("Pillow", "PIL.Image", "Pillow is needed for image recompression.")
        for page in pdf.pages:
            try:
                images = page.images
            except Exception:  # noqa: BLE001
                continue
            for _name, raw in list(images.items()):
                try:
                    pil = PdfImage(raw).as_pil_image()
                    if pil.mode in ("RGBA", "P", "LA", "CMYK"):
                        pil = pil.convert("RGB")
                    if max(pil.size) > max_px:
                        r = max_px / max(pil.size)
                        pil = pil.resize((max(1, int(pil.width * r)),
                                          max(1, int(pil.height * r))))
                    buf = io.BytesIO()
                    pil.save(buf, "JPEG", quality=int(quality), optimize=True)
                    raw.write(buf.getvalue(), filter=Name("/DCTDecode"))
                    raw.ColorSpace = Name("/DeviceGray") if pil.mode == "L" else Name("/DeviceRGB")
                    raw.BitsPerComponent = 8
                    if "/SMask" in raw:
                        del raw.SMask
                except Exception:  # noqa: BLE001
                    continue  # leave this image untouched
    bo = io.BytesIO()
    pdf.save(bo, compress_streams=True,
             object_stream_mode=pikepdf.ObjectStreamMode.generate)
    return bo.getvalue()


def _gs_compress(data, quality):
    # Map a 0-100 quality to Ghostscript presets.
    preset = "/screen" if quality < 45 else "/ebook" if quality < 75 else "/printer"
    with tempfile.TemporaryDirectory() as td:
        ip = os.path.join(td, "in.pdf")
        op = os.path.join(td, "out.pdf")
        with open(ip, "wb") as fh:
            fh.write(data)
        subprocess.run([gs_binary(), "-sDEVICE=pdfwrite",
                        "-dCompatibilityLevel=1.5", f"-dPDFSETTINGS={preset}",
                        "-dNOPAUSE", "-dQUIET", "-dBATCH",
                        f"-sOutputFile={op}", ip], check=True, timeout=300)
        with open(op, "rb") as fh:
            return fh.read()


# ---- conversions --------------------------------------------------------- #
_OFFICE_LOCK = threading.Lock()
OFFICE_EXTS = {"doc", "docx", "ppt", "pptx", "xls", "xlsx",
               "odt", "odp", "ods", "rtf", "txt", "csv"}
IMAGE_EXTS = {"png", "jpg", "jpeg", "bmp", "gif", "tif", "tiff", "webp"}


def conv_office2pdf(files, **_):
    soffice = office_binary()
    if not soffice:
        raise FeatureError(
            "LibreOffice is required for Office → PDF and was not found.\n"
            "Install it (free, all platforms): https://www.libreoffice.org/download\n"
            "macOS: brew install --cask libreoffice  •  Windows: winget install "
            "TheDocumentFoundation.LibreOffice")
    results = []
    for f in files:
        ext = os.path.splitext(f["name"])[1].lstrip(".").lower()
        if ext not in OFFICE_EXTS:
            raise FeatureError(f"“{f['name']}” isn't a supported Office file.")
        with _OFFICE_LOCK, tempfile.TemporaryDirectory() as td:
            ip = os.path.join(td, f"in.{ext}")
            with open(ip, "wb") as fh:
                fh.write(f["bytes"])
            profile = os.path.join(td, "profile")
            subprocess.run(
                [soffice, "--headless", "--norestore", "--nolockcheck",
                 f"-env:UserInstallation=file://{profile}",
                 "--convert-to", "pdf", "--outdir", td, ip],
                check=True, capture_output=True, timeout=300,
                env=dict(os.environ, HOME=td))
            outp = os.path.join(td, "in.pdf")
            if not os.path.exists(outp):
                raise FeatureError(f"LibreOffice could not convert “{f['name']}”.")
            with open(outp, "rb") as fh:
                results.append((f"{safe_stem(f['name'])}.pdf", fh.read()))
    return results, f"Converted {len(results)} file(s) to PDF."


def conv_pdf2docx(files, **_):
    need("pdf2docx")
    from pdf2docx import Converter
    results = []
    for f in files:
        with tempfile.TemporaryDirectory() as td:
            ip = os.path.join(td, "in.pdf")
            op = os.path.join(td, "out.docx")
            with open(ip, "wb") as fh:
                fh.write(f["bytes"])
            cv = Converter(ip)
            try:
                cv.convert(op)
            finally:
                cv.close()
            with open(op, "rb") as fh:
                results.append((f"{safe_stem(f['name'])}.docx", fh.read()))
    return results, f"Converted {len(results)} PDF(s) to Word."


def conv_pdf2pptx(files, dpi=150, **_):
    pdfium = need("pypdfium2")
    need("pptx", "pptx", "python-pptx is required.")
    from pptx import Presentation
    from pptx.util import Emu
    results = []
    scale = max(72, int(dpi)) / 72.0
    for f in files:
        doc = pdfium.PdfDocument(f["bytes"])
        prs = Presentation()
        w_pt, h_pt = doc[0].get_size()
        prs.slide_width = Emu(int(w_pt / 72 * 914400))
        prs.slide_height = Emu(int(h_pt / 72 * 914400))
        blank = prs.slide_layouts[6]
        for i in range(len(doc)):
            pil = doc[i].render(scale=scale).to_pil()
            ib = io.BytesIO()
            pil.save(ib, "PNG")
            ib.seek(0)
            slide = prs.slides.add_slide(blank)
            slide.shapes.add_picture(ib, 0, 0,
                                     width=prs.slide_width, height=prs.slide_height)
        bo = io.BytesIO()
        prs.save(bo)
        results.append((f"{safe_stem(f['name'])}.pptx", bo.getvalue()))
    return results, f"Converted {len(results)} PDF(s) to PowerPoint."


def conv_pdf2img(files, img_format="png", dpi=150, **_):
    pdfium = need("pypdfium2")
    fmt = "JPEG" if str(img_format).lower() in ("jpg", "jpeg") else "PNG"
    ext = "jpg" if fmt == "JPEG" else "png"
    scale = max(72, int(dpi)) / 72.0
    results = []
    for f in files:
        doc = pdfium.PdfDocument(f["bytes"])
        stem = safe_stem(f["name"])
        for i in range(len(doc)):
            pil = doc[i].render(scale=scale).to_pil()
            if fmt == "JPEG" and pil.mode != "RGB":
                pil = pil.convert("RGB")
            ib = io.BytesIO()
            if fmt == "JPEG":
                pil.save(ib, fmt, quality=90)
            else:
                pil.save(ib, fmt)
            results.append((f"{stem}_p{i + 1:03d}.{ext}", ib.getvalue()))
    return results, f"Rendered {len(results)} page image(s)."


def conv_img2pdf(files, **_):
    need("Pillow", "PIL.Image", "Pillow is required for Images → PDF.")
    from PIL import Image as PILImage
    pages = []
    for f in files:
        ext = os.path.splitext(f["name"])[1].lstrip(".").lower()
        if ext not in IMAGE_EXTS:
            continue
        im = PILImage.open(io.BytesIO(f["bytes"]))
        pages.append(im.convert("RGB"))
    if not pages:
        raise FeatureError("No supported image files were provided.")
    bo = io.BytesIO()
    pages[0].save(bo, "PDF", save_all=True, append_images=pages[1:])
    return [("images.pdf", bo.getvalue())], f"Combined {len(pages)} image(s) into one PDF."


CONVERT_ROUTES = {
    "office2pdf": conv_office2pdf,
    "pdf2docx": conv_pdf2docx,
    "pdf2pptx": conv_pdf2pptx,
    "pdf2img": conv_pdf2img,
    "img2pdf": conv_img2pdf,
}


# --------------------------------------------------------------------------- #
#  YouTube download (yt-dlp)
#
#  Downloads run in a background thread as a "job" the client polls, so the UI
#  can show real progress (%, speed, ETA) for multi-GB videos. Files land in a
#  per-job temp dir and are streamed straight from disk (never held in RAM),
#  then cleaned up after delivery.
# --------------------------------------------------------------------------- #
_YT_LOCK = threading.Lock()
_YT_JOBS: dict = {}  # id -> {status,pct,msg,speed,eta,name,path,dir,error}


def _yt_dlp():
    return need("yt-dlp", "yt_dlp",
                "Could not install 'yt-dlp'. Try: pip install yt-dlp")


def _yt_format(height, compat="best"):
    """Pick a format string for the requested max height + compatibility mode.

    Why this matters: YouTube stores everything above 1080p ONLY in VP9/AV1.
    QuickTime cannot decode VP9, so a "4K mp4" made by remuxing VP9 into an
    mp4 container won't open on macOS even though the download itself worked.
      best  → best streams regardless of codec (VP9/AV1; VLC/IINA territory)
      h264  → native-H.264 streams only, plays anywhere but capped ~1080p
      qt    → same as best; the worker then re-encodes to HEVC/H.264 mp4
    """
    h = f"[height<={int(height)}]" if height else ""
    if not ffmpeg_binary():
        return f"best{h}[ext=mp4]/best{h}/best"
    if compat == "h264":
        return (f"bestvideo{h}[vcodec^=avc1]+bestaudio[ext=m4a]/"
                f"best{h}[ext=mp4]/best[ext=mp4]/best")
    if compat == "qt":
        # Prefer AV1 over VP9: identical quality tier, but decodes far faster
        # (dav1d in software; hardware AV1 decode on Apple M3+/RTX 30+ with
        # ffmpeg 7.1+). 10-bit VP9 (HDR) software decode is the classic
        # bottleneck that pins conversion at ~1x realtime.
        return (f"bestvideo{h}[vcodec^=av01]+bestaudio/"
                f"bestvideo{h}+bestaudio/best{h}/best")
    return f"bestvideo{h}+bestaudio/best{h}/best"


_ENC_CACHE = None


def _enc_works(vargs):
    """ffmpeg lists hardware encoders even when the hardware/driver is absent
    (e.g. hevc_nvenc with no NVIDIA GPU) — so verify with a 1-frame encode."""
    try:
        r = subprocess.run(
            [ffmpeg_binary(), "-hide_banner", "-loglevel", "error",
             "-f", "lavfi", "-i", "color=black:size=640x360:rate=1",
             "-frames:v", "1", *vargs, "-f", "null", "-"],
            capture_output=True, timeout=30)
        return r.returncode == 0
    except Exception:  # noqa: BLE001
        return False


def _encoder_args():
    """Best available H.264/HEVC encoder: hardware if it actually works,
    x264 software fallback. Result cached for the process lifetime."""
    global _ENC_CACHE
    if _ENC_CACHE is not None:
        return _ENC_CACHE
    try:
        listed = subprocess.run([ffmpeg_binary(), "-hide_banner", "-encoders"],
                                capture_output=True, text=True, timeout=15).stdout
    except Exception:  # noqa: BLE001
        listed = ""
    candidates = []
    if sys.platform == "darwin":
        # hvc1 tag is REQUIRED for QuickTime to recognise HEVC in mp4.
        # -prio_speed tells the media engine to favor throughput (ffmpeg 5.1+);
        # plain variants follow as fallbacks for older ffmpeg builds.
        candidates += [
            ("hevc_videotoolbox",
             ["-c:v", "hevc_videotoolbox", "-q:v", "55", "-tag:v", "hvc1",
              "-prio_speed", "1"], "HEVC (hw)"),
            ("hevc_videotoolbox",
             ["-c:v", "hevc_videotoolbox", "-q:v", "55", "-tag:v", "hvc1"],
             "HEVC (hw)"),
            ("h264_videotoolbox",
             ["-c:v", "h264_videotoolbox", "-b:v", "14M", "-prio_speed", "1"],
             "H.264 (hw)"),
            ("h264_videotoolbox",
             ["-c:v", "h264_videotoolbox", "-b:v", "14M"], "H.264 (hw)"),
        ]
    candidates += [
        ("hevc_nvenc", ["-c:v", "hevc_nvenc", "-cq", "24", "-tag:v", "hvc1"],
         "HEVC (NVENC)"),
        ("h264_nvenc", ["-c:v", "h264_nvenc", "-cq", "23"], "H.264 (NVENC)"),
    ]
    for name, vargs, label in candidates:
        if name in listed and _enc_works(vargs):
            _ENC_CACHE = (vargs, label)
            return _ENC_CACHE
    _ENC_CACHE = (["-c:v", "libx264", "-crf", "20", "-preset", "fast",
                   "-pix_fmt", "yuv420p"], "H.264 (x264)")
    return _ENC_CACHE


def _decode_ladder(enc_label):
    """Decode configurations to try, fastest first.

    1. Full-GPU: decoded frames stay in GPU memory and go straight into the
       hardware encoder — zero CPU involvement, zero 4K frame copies. This is
       the difference between ~1x and several-x realtime for 4K60.
    2. HW decode with frames downloaded to CPU (works with any encoder).
    3. Plain software decode (always works).
    A config that can't handle the stream fails at startup, and we fall
    through to the next one."""
    lad = []
    if sys.platform == "darwin":
        if "(hw)" in (enc_label or ""):
            lad.append(["-hwaccel", "videotoolbox",
                        "-hwaccel_output_format", "videotoolbox"])
        lad.append(["-hwaccel", "videotoolbox"])
    elif "NVENC" in (enc_label or ""):
        lad.append(["-hwaccel", "cuda", "-hwaccel_output_format", "cuda"])
        lad.append(["-hwaccel", "cuda"])
    lad.append([])
    return lad


def _convert_qt(job, src, duration):
    """Re-encode to a QuickTime-friendly mp4, reporting progress via the job.
    Tries the fastest decode pipeline first, falling back if the stream isn't
    supported by a given hwaccel (fails within a second, so retries are cheap)."""
    venc, enc_name = _encoder_args()
    dst = os.path.splitext(src)[0] + ".qt.mp4"
    last_err = "unknown error"
    for dec in _decode_ladder(enc_name):
        args = [ffmpeg_binary(), "-y", *dec, "-i", src,
                *venc, "-c:a", "aac", "-b:a", "192k",
                "-threads", "0", "-movflags", "+faststart",
                "-progress", "pipe:1", "-nostats", "-loglevel", "error", dst]
        proc = subprocess.Popen(args, stdout=subprocess.PIPE,
                                stderr=subprocess.PIPE, text=True)
        with _YT_LOCK:
            job.update(msg=f"Converting for QuickTime · {enc_name}",
                       pct=0, speed=None, eta=None, rate=None)
        fps = ""
        for line in proc.stdout:  # ffmpeg -progress: out_time_*= is microseconds
            line = line.strip()
            if line.startswith("fps="):
                fps = line.split("=", 1)[1].strip()
            elif line.startswith("speed="):
                sp = line.split("=", 1)[1].strip()
                if sp and sp != "N/A":
                    with _YT_LOCK:
                        job["rate"] = sp + (f" · {fps} fps"
                                            if fps and fps != "0.00" else "")
            elif line.startswith(("out_time_ms=", "out_time_us=")) and duration:
                try:
                    us = int(line.split("=", 1)[1])
                except ValueError:
                    continue
                with _YT_LOCK:
                    job["pct"] = min(99.9, round(100 * (us / 1e6) / duration, 1))
        proc.wait()
        with _YT_LOCK:
            job["rate"] = None
        if proc.returncode == 0 and os.path.exists(dst) and os.path.getsize(dst):
            os.remove(src)
            # tidy name: "Title [id].qt.mp4" -> "Title [id].mp4"
            final = dst.replace(".qt.mp4", ".mp4")
            if final != dst and not os.path.exists(final):
                os.rename(dst, final)
                dst = final
            return dst, enc_name
        last_err = (proc.stderr.read() or "")[-500:] or last_err
        if os.path.exists(dst):
            os.remove(dst)  # clear partial output before the next attempt
    raise RuntimeError(f"ffmpeg conversion failed:\n{last_err}")


_VERTICAL_SIZES = {"1080": (1080, 1920), "720": (720, 1280), "1350": (1080, 1350)}


def _vertical_filters(mode, w, h):
    """Filter graphs that turn any aspect ratio into a vertical canvas.

    Returns a list to try in order — gblur is the nicer look but isn't in
    every ffmpeg build, so boxblur follows as a fallback.
    """
    fit = f"scale={w}:{h}:force_original_aspect_ratio=decrease"
    fill = f"scale={w}:{h}:force_original_aspect_ratio=increase,crop={w}:{h}"
    if mode == "crop":
        return [fill]
    if mode == "pad":
        return [f"{fit},pad={w}:{h}:(ow-iw)/2:(oh-ih)/2:black"]
    # "blur": the reel look — filled blurred backdrop, whole frame on top
    return [
        f"split[a][b];[a]{fill},gblur=sigma=28[bg];[b]{fit}[fg];"
        f"[bg][fg]overlay=(W-w)/2:(H-h)/2",
        f"split[a][b];[a]{fill},boxblur=20:2[bg];[b]{fit}[fg];"
        f"[bg][fg]overlay=(W-w)/2:(H-h)/2",
    ]


def _convert_vertical(job, src, duration, mode="blur", size="1080"):
    """Re-frame a video to a vertical canvas, reporting progress via the job."""
    venc, enc_name = _encoder_args()
    w, h = _VERTICAL_SIZES.get(str(size), (1080, 1920))
    dst = os.path.splitext(src)[0] + f".vertical.mp4"
    label = {"blur": "blurred backdrop", "crop": "cropped to fill",
             "pad": "black bars"}.get(mode, mode)
    last_err = "unknown error"

    for vf in _vertical_filters(mode, w, h):
        args = [ffmpeg_binary(), "-y", "-i", src,
                "-filter_complex" if "[" in vf else "-vf", vf,
                *venc, "-c:a", "aac", "-b:a", "192k",
                "-r", "30", "-threads", "0", "-movflags", "+faststart",
                "-progress", "pipe:1", "-nostats", "-loglevel", "error", dst]
        proc = subprocess.Popen(args, stdout=subprocess.PIPE,
                                stderr=subprocess.PIPE, text=True)
        with _YT_LOCK:
            job.update(msg=f"Reframing to {w}x{h} · {label}",
                       pct=0, speed=None, eta=None, rate=None)
        fps = ""
        for line in proc.stdout:
            line = line.strip()
            if line.startswith("fps="):
                fps = line.split("=", 1)[1].strip()
            elif line.startswith("speed="):
                sp = line.split("=", 1)[1].strip()
                if sp and sp != "N/A":
                    with _YT_LOCK:
                        job["rate"] = sp + (f" · {fps} fps"
                                            if fps and fps != "0.00" else "")
            elif line.startswith(("out_time_ms=", "out_time_us=")) and duration:
                try:
                    us = int(line.split("=", 1)[1])
                except ValueError:
                    continue
                with _YT_LOCK:
                    job["pct"] = min(99.9, round(100 * (us / 1e6) / duration, 1))
        proc.wait()
        with _YT_LOCK:
            job["rate"] = None
        if proc.returncode == 0 and os.path.exists(dst) and os.path.getsize(dst):
            os.remove(src)
            return dst, f"{w}x{h} {label} · {enc_name}"
        last_err = (proc.stderr.read() or "")[-500:] or last_err
        if os.path.exists(dst):
            os.remove(dst)
    raise RuntimeError(f"Vertical conversion failed:\n{last_err}")


def op_yt_info(url, cookies=""):
    yt = _yt_dlp()
    opts = {"quiet": True, "no_warnings": True, "noplaylist": True,
            "skip_download": True}
    if cookies:
        opts["cookiesfrombrowser"] = (cookies,)
    try:
        with yt.YoutubeDL(opts) as ydl:
            info = ydl.extract_info(url, download=False)
    except Exception as e:  # noqa: BLE001
        raise FeatureError(f"Could not read that URL:\n{e}")
    if info.get("_type") == "playlist":  # noplaylist usually handles this
        entries = info.get("entries") or []
        info = entries[0] if entries else info
    heights = sorted({f.get("height") for f in info.get("formats", [])
                      if f.get("height") and f.get("vcodec") not in (None, "none")},
                     reverse=True)
    return {
        "title": info.get("title") or "video",
        "uploader": info.get("uploader") or info.get("channel") or "",
        "duration": info.get("duration") or 0,
        "thumbnail": info.get("thumbnail") or "",
        "heights": heights[:12],
        "has_ffmpeg": bool(ffmpeg_binary()),
    }


def _yt_worker(job_id, url, height, compat, vertical="off",
               vsize="1080", cookies=""):
    yt = _yt_dlp()
    job = _YT_JOBS[job_id]

    def hook(d):
        with _YT_LOCK:
            if d.get("status") == "downloading":
                tot = d.get("total_bytes") or d.get("total_bytes_estimate") or 0
                got = d.get("downloaded_bytes") or 0
                job["pct"] = round(100.0 * got / tot, 1) if tot else None
                job["speed"] = d.get("speed")
                job["eta"] = d.get("eta")
                job["msg"] = "Downloading"
            elif d.get("status") == "finished":
                job["pct"] = 100
                job["speed"] = None
                job["eta"] = None
                job["msg"] = "Processing (merging streams)"

    opts = {
        "quiet": True, "no_warnings": True, "noplaylist": True,
        # .120B = truncate the (possibly unicode) title to 120 bytes safely
        "outtmpl": os.path.join(job["dir"], "%(title).120B [%(id)s].%(ext)s"),
        "format": _yt_format(height, compat),
        "progress_hooks": [hook],
        "retries": 3,
    }
    if cookies:
        # Instagram/private posts need a logged-in session; yt-dlp can borrow
        # the cookies straight out of an installed browser.
        opts["cookiesfrombrowser"] = (cookies,)
    if compat == "h264" and ffmpeg_binary():
        opts["merge_output_format"] = "mp4"
    # compat 'best'/'qt': let yt-dlp pick a container that fits the codecs
    # (webm/mkv) — forcing VP9 into mp4 is exactly what confuses QuickTime.
    try:
        with yt.YoutubeDL(opts) as ydl:
            info = ydl.extract_info(url, download=True)

        # Final file: yt-dlp reports it directly; fall back to biggest file.
        path = None
        for rd in (info or {}).get("requested_downloads") or []:
            if rd.get("filepath") and os.path.isfile(rd["filepath"]):
                path = rd["filepath"]
                break
        if not path:
            paths = [os.path.join(job["dir"], p) for p in os.listdir(job["dir"])]
            paths = [p for p in paths if os.path.isfile(p)
                     and not p.endswith((".part", ".ytdl", ".temp"))]
            if not paths:
                raise RuntimeError("yt-dlp finished but produced no file.")
            path = max(paths, key=os.path.getsize)

        # What did we actually get? (verifiable answer to "was it really 4K?")
        vfmt = ((info or {}).get("requested_formats") or [info or {}])[0]
        w, h = vfmt.get("width"), vfmt.get("height")
        vc = (vfmt.get("vcodec") or "?").split(".")[0]
        detail = f"{w}×{h} · {vc}" if w and h else vc

        dur = (info or {}).get("duration")
        if vertical and vertical != "off":
            if not ffmpeg_binary():
                raise RuntimeError("Vertical reframing needs ffmpeg. Install it "
                                   "and restart, or set Vertical to “Keep original”.")
            path, note = _convert_vertical(job, path, dur, vertical, vsize)
            detail += f" → {note}"
        elif compat == "qt" and ffmpeg_binary():
            path, enc = _convert_qt(job, path, dur)
            detail += f" → {enc}"

        with _YT_LOCK:
            job.update(status="done", pct=100, msg="Done", detail=detail,
                       path=path, name=os.path.basename(path))
    except Exception as e:  # noqa: BLE001
        with _YT_LOCK:
            job.update(status="error", error=str(e), msg="Failed")


def op_yt_start(url, height, compat="best", vertical="off",
                vsize="1080", cookies=""):
    if not url or not re.match(r"https?://", url.strip()):
        raise FeatureError("That doesn't look like a URL. Paste the full "
                           "https://… video link.")
    if compat not in ("best", "h264", "qt"):
        compat = "best"
    if vertical not in ("off", "blur", "crop", "pad"):
        vertical = "off"
    if cookies not in ("", "chrome", "firefox", "safari", "edge", "brave",
                       "chromium", "opera", "vivaldi"):
        cookies = ""
    job_id = uuid.uuid4().hex
    _YT_JOBS[job_id] = {"status": "running", "pct": None, "msg": "Starting",
                        "speed": None, "eta": None, "name": None, "path": None,
                        "detail": None, "rate": None, "error": None,
                        "dir": tempfile.mkdtemp(prefix="ytdl_")}
    threading.Thread(target=_yt_worker,
                     args=(job_id, url.strip(), height, compat, vertical,
                           vsize, cookies), daemon=True).start()
    return {"id": job_id}


def _yt_cleanup(job_id):
    with _YT_LOCK:
        job = _YT_JOBS.pop(job_id, None)
    if job and job.get("dir"):
        shutil.rmtree(job["dir"], ignore_errors=True)


# --------------------------------------------------------------------------- #
#  QR codes
# --------------------------------------------------------------------------- #
_EC_LEVELS = {"L": 1, "M": 0, "Q": 3, "H": 2}   # qrcode.constants values
_QR_MAX_PX = 4096


def _hex_rgba(value, fallback=(0, 0, 0, 255)):
    """'#rgb' / '#rrggbb' / '#rrggbbaa' / 'transparent' -> (r,g,b,a)."""
    s = str(value or "").strip().lower()
    if s in ("transparent", "none", ""):
        return (0, 0, 0, 0)
    s = s.lstrip("#")
    if len(s) == 3:
        s = "".join(c * 2 for c in s)
    if len(s) == 6:
        s += "ff"
    if len(s) != 8 or not re.fullmatch(r"[0-9a-f]{8}", s):
        return fallback
    return tuple(int(s[i:i + 2], 16) for i in (0, 2, 4, 6))


def _relative_luminance(rgba):
    def chan(c):
        c /= 255.0
        return c / 12.92 if c <= 0.03928 else ((c + 0.055) / 1.055) ** 2.4
    r, g, b = rgba[0], rgba[1], rgba[2]
    return 0.2126 * chan(r) + 0.7152 * chan(g) + 0.0722 * chan(b)


def _contrast_ratio(a, b):
    la, lb = _relative_luminance(a), _relative_luminance(b)
    hi, lo = max(la, lb), min(la, lb)
    return (hi + 0.05) / (lo + 0.05)


def qr_payload(kind, fields):
    """Build the string that actually goes into the QR code."""
    f = {k: str(v or "").strip() for k, v in (fields or {}).items()}
    kind = (kind or "text").lower()

    def esc(s):  # Wi-Fi / vCard escaping
        return re.sub(r"([\\;,:\"])", r"\\\1", s)

    if kind == "url":
        u = f.get("text", "")
        if u and not re.match(r"^[a-z][a-z0-9+.\-]*:", u, re.I):
            u = "https://" + u
        return u
    if kind == "wifi":
        enc = (f.get("security") or "WPA").upper()
        if enc == "NONE":
            return f"WIFI:T:nopass;S:{esc(f.get('ssid',''))};;"
        return ("WIFI:T:{t};S:{s};P:{p};{h}".format(
            t=enc, s=esc(f.get("ssid", "")), p=esc(f.get("password", "")),
            h="H:true;;" if f.get("hidden") in ("1", "true", "True") else ";"))
    if kind == "email":
        q = []
        if f.get("subject"):
            q.append("subject=" + _pct(f["subject"]))
        if f.get("body"):
            q.append("body=" + _pct(f["body"]))
        return "mailto:" + f.get("to", "") + ("?" + "&".join(q) if q else "")
    if kind == "sms":
        return "SMSTO:{n}:{m}".format(n=f.get("phone", ""), m=f.get("message", ""))
    if kind == "phone":
        return "tel:" + f.get("phone", "")
    if kind == "geo":
        return "geo:{la},{lo}".format(la=f.get("lat", "0"), lo=f.get("lon", "0"))
    if kind == "vcard":
        name = f.get("name", "")
        given, _, family = name.partition(" ")   # vCard N is Family;Given;...
        lines = ["BEGIN:VCARD", "VERSION:3.0",
                 f"N:{esc(family or name)};{esc(given if family else '')};;;",
                 f"FN:{esc(name)}"]
        if f.get("org"):
            lines.append("ORG:" + esc(f["org"]))
        if f.get("title"):
            lines.append("TITLE:" + esc(f["title"]))
        if f.get("phone"):
            lines.append("TEL;TYPE=CELL:" + f["phone"])
        if f.get("email"):
            lines.append("EMAIL:" + f["email"])
        if f.get("url"):
            lines.append("URL:" + f["url"])
        lines.append("END:VCARD")
        return "\n".join(lines)
    return f.get("text", "")


def _pct(s):
    from urllib.parse import quote
    return quote(s, safe="")


def _silhouette(img, colour):
    """Flatten any logo to a single-colour silhouette (the look in the sample)."""
    from PIL import Image
    img = img.convert("RGBA")
    alpha = img.getchannel("A")
    # If the image is essentially opaque, derive the shape from darkness instead.
    if alpha.getextrema()[0] > 250:
        grey = img.convert("L")
        mask = grey.point(lambda v: 255 if v < 128 else 0).convert("L")
    else:
        mask = alpha.point(lambda v: 255 if v > 40 else 0).convert("L")
    flat = Image.new("RGBA", img.size, tuple(colour))
    flat.putalpha(mask)
    return flat


def op_qr(data="", ec="H", target_px=1024, border=4, fg="#000000", bg="#ffffff",
          style="square", logo_data=None, logo_pct=22, logo_style="original",
          pad=True, pad_shape="rounded", pad_pct=6, fmt="png"):
    """Render a QR code, optionally with a centred logo. Returns a dict."""
    qrcode = need("qrcode")
    need("Pillow", "PIL")
    from PIL import Image, ImageDraw

    data = data if isinstance(data, str) else str(data or "")
    if not data.strip():
        raise FeatureError("Type something to encode first.")

    ec = (ec or "H").upper()
    if ec not in _EC_LEVELS:
        ec = "H"
    if logo_data and ec != "H":
        ec = "H"          # a logo eats modules; H (~30% recovery) is the floor
    border = max(0, min(16, int(border)))
    logo_pct = max(0, min(40, float(logo_pct)))

    try:
        q = qrcode.QRCode(error_correction=_EC_LEVELS[ec], box_size=1,
                          border=border)
        q.add_data(data)
        q.make(fit=True)
    except Exception as e:                                    # noqa: BLE001
        raise FeatureError(f"Could not encode that: {e}")

    matrix = q.get_matrix()
    n = len(matrix)                       # side length in modules, incl. border
    mods = n - 2 * border                 # the code itself

    target_px = max(128, min(_QR_MAX_PX, int(target_px or 1024)))
    box = max(1, round(target_px / n))
    size = n * box

    fg_rgba, bg_rgba = _hex_rgba(fg, (0, 0, 0, 255)), _hex_rgba(bg, (255, 255, 255, 255))
    img = Image.new("RGBA", (size, size), bg_rgba)
    draw = ImageDraw.Draw(img)

    def is_finder(mx, my):
        return ((mx < 7 and my < 7)
                or (mx >= mods - 7 and my < 7)
                or (mx < 7 and my >= mods - 7))

    style = (style or "square").lower()
    radius = max(1, int(box * 0.32))
    for y, row in enumerate(matrix):
        for x, on in enumerate(row):
            if not on:
                continue
            x0, y0 = x * box, y * box
            x1, y1 = x0 + box - 1, y0 + box - 1
            # Finder patterns always stay crisp squares — scanners lock onto them.
            if style == "square" or is_finder(x - border, y - border):
                draw.rectangle([x0, y0, x1, y1], fill=fg_rgba)
            elif style == "dots":
                draw.ellipse([x0, y0, x1, y1], fill=fg_rgba)
            else:  # rounded
                draw.rounded_rectangle([x0, y0, x1, y1], radius=radius, fill=fg_rgba)

    logo_note = ""
    if logo_data and logo_pct > 0:
        try:
            raw = base64.b64decode(logo_data.split(",")[-1])
            logo = Image.open(io.BytesIO(raw))
            logo.load()
        except Exception:                                     # noqa: BLE001
            raise FeatureError("That logo file couldn't be read as an image.")
        logo = logo.convert("RGBA")
        if (logo_style or "original") == "silhouette":
            logo = _silhouette(logo, fg_rgba)

        code_px = mods * box                         # drawable area, sans border
        want = max(1, int(code_px * logo_pct / 100.0))
        logo.thumbnail((want, want), Image.LANCZOS)
        lw, lh = logo.size
        cx, cy = size // 2, size // 2
        px = max(0, int(max(lw, lh) * max(0, min(50, float(pad_pct))) / 100.0))

        if pad:
            pw, ph = lw + 2 * px, lh + 2 * px
            bx = [cx - pw // 2, cy - ph // 2, cx + pw // 2, cy + ph // 2]
            pad_fill = bg_rgba if bg_rgba[3] else (255, 255, 255, 255)
            if pad_shape == "circle":
                d = max(pw, ph)
                draw.ellipse([cx - d // 2, cy - d // 2, cx + d // 2, cy + d // 2],
                             fill=pad_fill)
            elif pad_shape == "square":
                draw.rectangle(bx, fill=pad_fill)
            else:
                draw.rounded_rectangle(bx, radius=max(2, int(min(pw, ph) * 0.18)),
                                       fill=pad_fill)

        img.alpha_composite(logo, (cx - lw // 2, cy - lh // 2))
        covered = (lw * lh) / float(code_px * code_px) * 100
        if covered > 25:
            logo_note = (f"Logo covers ~{covered:.0f}% of the code — above ~25% "
                         "some scanners start to struggle. Test before you print.")

    out = io.BytesIO()
    if fmt == "jpg":
        flat = Image.new("RGB", img.size,
                         bg_rgba[:3] if bg_rgba[3] else (255, 255, 255))
        flat.paste(img, mask=img.getchannel("A"))
        flat.save(out, "JPEG", quality=95)
        mime = "image/jpeg"
    else:
        img.save(out, "PNG", optimize=True)
        mime = "image/png"
    png_b64 = base64.b64encode(out.getvalue()).decode()

    warn = []
    if logo_note:
        warn.append(logo_note)
    if border < 4:
        warn.append(f"Quiet zone is {border} modules — the spec wants 4. Below "
                    "that, scanners lose the edge against a busy background.")
    if not bg_rgba[3]:
        warn.append("Transparent background: only scannable on a light surface. "
                    "On anything dark it becomes invisible.")
    if bg_rgba[3] and _contrast_ratio(fg_rgba, bg_rgba) < 3.0:
        warn.append("Foreground and background are too close in brightness — "
                    "many scanners will fail. Darken the foreground.")
    elif _relative_luminance(fg_rgba) > _relative_luminance(bg_rgba):
        warn.append("Light code on a dark background — some older scanners "
                    "only read dark-on-light. Test it.")
    if len(data) > 1200:
        warn.append("That's a lot of data — the code is dense and needs a "
                    "clean, large print to scan reliably.")

    return {"image": png_b64, "mime": mime, "size": size, "modules": mods,
            "version": q.version, "ec": ec, "bytes": len(out.getvalue()),
            "chars": len(data), "warn": " ".join(warn)}


# --------------------------------------------------------------------------- #
#  Image fetch / crop / resize
# --------------------------------------------------------------------------- #
_IMG_PREVIEW_PX = 1400
_IMG_MAX_BYTES = 80 * 1024 * 1024


def _open_image(raw, label="image"):
    need("Pillow", "PIL")
    from PIL import Image, ImageOps
    try:
        img = Image.open(io.BytesIO(raw))
        img.load()
    except Exception:                                         # noqa: BLE001
        raise FeatureError(f"That {label} couldn't be read as an image.")
    return ImageOps.exif_transpose(img)


def img_describe(raw, name="image"):
    """Preview + dimensions for the crop UI. Original bytes stay server-side."""
    need("Pillow", "PIL")
    from PIL import Image
    img = _open_image(raw, name)
    w, h = img.size
    fmt = (img.format or "PNG").upper()
    prev = img.convert("RGB") if img.mode not in ("RGB", "L") else img
    prev = prev.copy()
    prev.thumbnail((_IMG_PREVIEW_PX, _IMG_PREVIEW_PX), Image.LANCZOS)
    buf = io.BytesIO()
    prev.convert("RGB").save(buf, "JPEG", quality=86)
    return {"name": name, "width": w, "height": h, "format": fmt,
            "bytes": len(raw), "preview": base64.b64encode(buf.getvalue()).decode()}


def op_img_process(raw, name="image", crop=None, out_w=None, out_h=None,
                   fmt="png", quality=90, rotate=0, flip=False, grayscale=False):
    need("Pillow", "PIL")
    from PIL import Image
    img = _open_image(raw, name)
    w, h = img.size

    if crop:
        x = max(0, min(w - 1, int(crop.get("x", 0))))
        y = max(0, min(h - 1, int(crop.get("y", 0))))
        cw = max(1, min(w - x, int(crop.get("w", w))))
        ch = max(1, min(h - y, int(crop.get("h", h))))
        img = img.crop((x, y, x + cw, y + ch))

    rotate = int(rotate or 0) % 360
    if rotate in (90, 180, 270):
        img = img.rotate(-rotate, expand=True)
    if flip:
        img = img.transpose(Image.FLIP_LEFT_RIGHT)
    if grayscale:
        img = img.convert("L").convert("RGBA" if fmt == "png" else "RGB")

    if out_w or out_h:
        cw, ch = img.size
        tw = int(out_w) if out_w else max(1, round(cw * int(out_h) / ch))
        th = int(out_h) if out_h else max(1, round(ch * int(out_w) / cw))
        tw, th = max(1, min(20000, tw)), max(1, min(20000, th))
        img = img.resize((tw, th), Image.LANCZOS)

    stem = re.sub(r"[^\w.\- ]+", "_", (name or "image").rsplit(".", 1)[0]) or "image"
    buf = io.BytesIO()
    fmt = (fmt or "png").lower()
    if fmt in ("jpg", "jpeg"):
        bg = Image.new("RGB", img.size, (255, 255, 255))
        rgba = img.convert("RGBA")
        bg.paste(rgba, mask=rgba.getchannel("A"))
        bg.save(buf, "JPEG", quality=int(quality), optimize=True, progressive=True)
        ext = "jpg"
    elif fmt == "webp":
        img.convert("RGBA").save(buf, "WEBP", quality=int(quality), method=5)
        ext = "webp"
    else:
        img.convert("RGBA").save(buf, "PNG", optimize=True)
        ext = "png"
    data = buf.getvalue()
    return ([(f"{stem}_{img.size[0]}x{img.size[1]}.{ext}", data)],
            f"{img.size[0]}×{img.size[1]} · {len(data)/1024:.0f} KB")


def op_img_fetch(url):
    """Download an image by URL so it can be cropped/resized locally."""
    from urllib.request import Request, urlopen
    url = (url or "").strip()
    if not re.match(r"^https?://", url, re.I):
        raise FeatureError("Paste a full image URL starting with http:// or https://")
    req = Request(url, headers={
        "User-Agent": ("Mozilla/5.0 (Macintosh; Intel Mac OS X 10_15_7) "
                       "AppleWebKit/537.36 (KHTML, like Gecko) "
                       "Chrome/125.0 Safari/537.36"),
        "Accept": "image/avif,image/webp,image/apng,image/*,*/*;q=0.8",
    })
    try:
        with urlopen(req, timeout=30) as r:
            ctype = (r.headers.get("Content-Type") or "").split(";")[0].strip()
            raw = r.read(_IMG_MAX_BYTES + 1)
    except Exception as e:  # noqa: BLE001
        raise FeatureError(f"Could not fetch that URL:\n{e}")
    if len(raw) > _IMG_MAX_BYTES:
        raise FeatureError("That image is over 80 MB — save it locally and drop "
                           "it in instead.")
    if ctype and not ctype.startswith("image/") and "octet-stream" not in ctype:
        raise FeatureError(f"That URL returned {ctype or 'no image'}, not an "
                           "image. Use the direct link to the image file "
                           "(right-click → Copy image address).")
    name = os.path.basename(urlparse(url).path) or "image"
    if "." not in name:
        name += "." + (ctype.split("/")[-1] if ctype else "jpg").replace("jpeg", "jpg")
    return raw, name


# --------------------------------------------------------------------------- #
#  PDF -> Markdown  (token-lean output for LLMs, with figures kept accurate)
#
#  The goal is not "a picture of the paper". It is: text stays text (cheap),
#  figures become cropped images anchored exactly where they occur in the
#  reading order, and anything we cannot read confidently is escalated to a
#  page image rather than silently mangled.
# --------------------------------------------------------------------------- #

fitz = None                 # bound by _md_init(); helpers are testable alone


def _md_init():
    """Import PyMuPDF once and publish it to this module's globals."""
    global fitz
    if fitz is None:
        fitz = need("PyMuPDF", "fitz")
    return fitz


_MD_BODY_TOL = 0.6          # font pt within this of the mode counts as body text
_MD_FURNITURE_HITS = 0.6    # a line repeating on >=60% of pages is header/footer
_MD_FIG_GAP = 14            # pt: drawings closer than this merge into one figure
_MD_FIG_MIN_AREA = 0.012    # ignore vector specks below this fraction of a page
_MD_CAP_GAP = 74            # pt: how far from a figure a caption may sit
_MD_MATH = "∑∏∫√≤≥≠≈±×÷αβγδθλμσφω∂∇∈∉⊂⊆∪∩→⇒↔∞"

_CAP_RE = re.compile(
    r"^\s*(Figure|Fig\.?|Table|Chart|Algorithm|Listing|Scheme|Exhibit)\s*"
    r"([0-9]+(?:\.[0-9]+)?|[IVXLC]+|[A-Z])\s*[.:\u2014-]?\s*", re.I)
_NUM_HEAD_RE = re.compile(
    r"^\s*((?:[0-9]+\.){0,3}[0-9]+|[IVXLC]+\.|Appendix\s+[A-Z])\s+\S")
_WORD_HEAD_RE = re.compile(
    r"^\s*(abstract|introduction|background|related work|method(?:s|ology)?|"
    r"approach|experiment(?:s|al setup)?|evaluation|results?|discussion|"
    r"conclusions?|future work|references|bibliography|acknowledg(?:e)?ments?|"
    r"appendix)\s*$", re.I)
_FOOT_RE = re.compile(r"^\s*[\u00b9\u00b2\u00b3\u2070-\u209f\*\u2020\u2021]|^\s*\[?\d{1,2}[\].]\s")


def _md_spans(page):
    """Flat list of text spans with geometry. One dict per span, reading-agnostic."""
    out = []
    d = page.get_text("dict")
    for blk in d.get("blocks", []):
        if blk.get("type") != 0:
            continue
        for line in blk.get("lines", []):
            for sp in line.get("spans", []):
                t = sp.get("text", "")
                if not t.strip():
                    continue
                out.append({
                    "text": t, "size": round(sp.get("size", 0), 1),
                    "font": sp.get("font", ""), "flags": sp.get("flags", 0),
                    "bbox": fitz.Rect(sp["bbox"]),
                    "line": tuple(round(v, 1) for v in line["bbox"]),
                })
    return out


def _md_lines(spans):
    """Group spans back into visual lines, preserving their order on the line."""
    by = {}
    for sp in spans:
        by.setdefault(sp["line"], []).append(sp)
    lines = []
    for key, group in by.items():
        group.sort(key=lambda s: s["bbox"].x0)
        r = fitz.Rect(group[0]["bbox"])
        for s in group[1:]:
            r |= s["bbox"]
        text = "".join(s["text"] for s in group)
        sizes = {}
        for s in group:
            sizes[s["size"]] = sizes.get(s["size"], 0) + len(s["text"])
        lines.append({
            "text": _md_tidy(text), "bbox": r,
            "size": max(sizes, key=sizes.get) if sizes else 0,
            "bold": any("bold" in s["font"].lower() or (s["flags"] & 16)
                        for s in group),
            "italic": any("italic" in s["font"].lower() or "it" in s["font"].lower()
                          or (s["flags"] & 2) for s in group),
            "spans": group,
        })
    lines.sort(key=lambda l: (round(l["bbox"].y0, 1), l["bbox"].x0))
    return lines


def _md_tidy(s):
    s = s.replace("\ufb00", "ff").replace("\ufb01", "fi").replace("\ufb02", "fl")
    s = s.replace("\ufb03", "ffi").replace("\ufb04", "ffl")
    s = s.replace("\u00ad", "").replace("\u2010", "-")
    return re.sub(r"[ \t]+", " ", s).strip()


def _md_body_size(pages_lines):
    """The most common font size, weighted by characters — that is body text."""
    tally = {}
    for lines in pages_lines:
        for l in lines:
            tally[l["size"]] = tally.get(l["size"], 0) + len(l["text"])
    return max(tally, key=tally.get) if tally else 10.0


def _md_furniture(pages_lines, page_rects, body_size=None):
    """Running headers/footers: same small text, same band, on most pages.

    Size matters: a paper's running head is often the title verbatim, so
    without the size guard the real title gets deleted from page one.
    """
    if len(pages_lines) < 3:
        return set()
    seen = {}
    for lines, rect in zip(pages_lines, page_rects):
        top, bot = rect.height * 0.08, rect.height * 0.92
        for l in lines:
            if body_size and l["size"] > body_size + 1.0:
                continue
            if l["bbox"].y1 <= top or l["bbox"].y0 >= bot:
                key = re.sub(r"\d+", "#", l["text"])[:70]
                if key:
                    seen.setdefault(key, set()).add(id(lines))
    need = max(2, int(len(pages_lines) * _MD_FURNITURE_HITS))
    return {k for k, v in seen.items() if len(v) >= need}


def _md_columns(lines, rect):
    """Find column gutters by projecting text onto the x-axis and looking for
    vertical bands that no ordinary line occupies.

    Generalises to any column count: journals use two, some proceedings use
    three, and most documents use one. A few things legitimately span the
    gutters — the title, a full-width figure — so those are excluded from the
    projection rather than treated as counter-evidence.
    """
    if len(lines) < 8:
        return []
    wide = rect.width * 0.55
    body = [l for l in lines if l["bbox"].width < wide]
    if len(body) < 6:
        return []

    step = 3.0
    nbin = max(8, int(rect.width / step) + 1)
    hits = [0] * nbin
    for l in body:
        a = max(0, int((l["bbox"].x0 - rect.x0) / step))
        b = min(nbin - 1, int((l["bbox"].x1 - rect.x0) / step))
        for i in range(a, b + 1):
            hits[i] += 1

    # Ignore the outer margins; we only care about gaps between text blocks.
    used = [i for i, h in enumerate(hits) if h]
    if not used:
        return []
    lo, hi = used[0], used[-1]

    gaps, run = [], None
    for i in range(lo, hi + 1):
        if hits[i] == 0:
            run = i if run is None else run
        elif run is not None:
            gaps.append((run, i - 1))
            run = None
    min_gap = max(8.0, rect.width * 0.018)
    cuts = []
    for a, b in gaps:
        if (b - a + 1) * step < min_gap:
            continue
        x = rect.x0 + (a + b + 1) / 2 * step
        left = sum(1 for l in body if l["bbox"].x1 <= x)
        right = sum(1 for l in body if l["bbox"].x0 >= x)
        if left >= 3 and right >= 3:
            cuts.append(x)
    if len(cuts) > 4:                       # implausible; treat as one column
        return []
    return sorted(cuts)


def _md_order(lines, rect, cuts):
    """Reading order. With columns, read each column fully before moving to the
    next; anything spanning a gutter (title, wide figure, full-width table) acts
    as a horizontal barrier that resets the column sweep beneath it."""
    if not cuts:
        return sorted(lines, key=lambda l: (round(l["bbox"].y0, 1), l["bbox"].x0))

    edges = [rect.x0 - 1] + list(cuts) + [rect.x1 + 1]

    def col_of(l):
        """Which column a line belongs to, or None if it spans a gutter."""
        for c in cuts:
            if l["bbox"].x0 < c - 6 and l["bbox"].x1 > c + 6:
                return None
        mid = (l["bbox"].x0 + l["bbox"].x1) / 2
        for i in range(len(edges) - 1):
            if edges[i] <= mid < edges[i + 1]:
                return i
        return 0

    spanning = [l for l in lines if col_of(l) is None]
    bands, prev = [], rect.y0 - 1
    for l in sorted(spanning, key=lambda l: l["bbox"].y0):
        if l["bbox"].y0 - prev > 4:
            bands.append((prev, l["bbox"].y0))
        bands.append((l["bbox"].y0, l["bbox"].y1))
        prev = l["bbox"].y1
    bands.append((prev, rect.y1 + 1))

    out, placed = [], set()
    for y0, y1 in bands:
        here = [l for l in lines
                if y0 <= l["bbox"].y0 < y1 and id(l) not in placed]
        if not here:
            continue
        span_here = [l for l in here if col_of(l) is None]
        if span_here:
            for l in sorted(span_here, key=lambda l: l["bbox"].y0):
                out.append(l)
                placed.add(id(l))
            continue
        for ci in range(len(edges) - 1):
            col = [l for l in here if col_of(l) == ci]
            for l in sorted(col, key=lambda l: (round(l["bbox"].y0, 1),
                                                l["bbox"].x0)):
                out.append(l)
                placed.add(id(l))
    for l in lines:                      # nothing may be dropped
        if id(l) not in placed:
            out.append(l)
    return out


def _md_cluster_rects(rects, gap):
    """Merge rectangles that touch or nearly touch into connected groups."""
    boxes = [fitz.Rect(r) for r in rects]
    merged = True
    while merged and boxes:
        merged = False
        for i in range(len(boxes)):
            for j in range(i + 1, len(boxes)):
                a, b = boxes[i], boxes[j]
                grown = fitz.Rect(a.x0 - gap, a.y0 - gap, a.x1 + gap, a.y1 + gap)
                if grown.intersects(b):
                    boxes[i] = a | b
                    boxes.pop(j)
                    merged = True
                    break
            if merged:
                break
    return boxes


def _md_table_regions(page):
    """Table regions from ruling geometry: horizontal rules or a cell grid."""
    hlines, cells = [], []
    try:
        drawings = page.get_drawings()
    except Exception:
        return []
    for d in drawings:
        r = fitz.Rect(d["rect"])
        for item in d.get("items", []):
            if item[0] == "l":
                if abs(item[1].y - item[2].y) < 1.2 and abs(item[1].x - item[2].x) > 28:
                    hlines.append(fitz.Rect(min(item[1].x, item[2].x), item[1].y - .6,
                                            max(item[1].x, item[2].x), item[1].y + .6))
            elif item[0] == "re":
                rr = fitz.Rect(item[1])
                if rr.width > 18 and 6 < rr.height < 46:
                    cells.append(rr)
                elif rr.height < 1.5 and rr.width > 28:
                    hlines.append(rr)

    regions = []
    if len(cells) >= 4:                       # full cell grid (Word-style)
        for box in _md_cluster_rects(cells, 3):
            inside = [c for c in cells if box.intersects(c)]
            if len(inside) >= 4:
                regions.append(box)
    if len(hlines) >= 2:                      # booktabs-style rules
        for box in _md_cluster_rects(hlines, 46):
            rules = [h for h in hlines if box.intersects(h)]
            if len(rules) >= 2 and box.height > 12:
                if not any(box.intersects(r) and (box & r).get_area() >
                           0.5 * box.get_area() for r in regions):
                    regions.append(box)
    return regions


def _md_table_from(region, lines):
    """Cluster the text inside a region into rows and columns.

    Works for both ruled grids and horizontal-rule-only tables, because it
    keys off where the text actually sits rather than the ruling style.
    """
    inner = [l for l in lines
             if region.intersects(l["bbox"])
             and (region & l["bbox"]).get_area() > 0.45 * l["bbox"].get_area()]
    if len(inner) < 4:
        return None

    rows = []                                  # cluster by vertical overlap
    for l in sorted(inner, key=lambda l: l["bbox"].y0):
        placed = False
        for row in rows:
            if abs(row[0]["bbox"].y0 - l["bbox"].y0) < max(4, l["bbox"].height * 0.6):
                row.append(l)
                placed = True
                break
        if not placed:
            rows.append([l])
    rows = [sorted(r, key=lambda l: l["bbox"].x0) for r in rows]
    rows = [r for r in rows if r]
    if len(rows) < 2:
        return None

    # Column edges: the set of distinct left-edges across all rows.
    edges = []
    for r in rows:
        for c in r:
            x = c["bbox"].x0
            if not any(abs(x - e) < 12 for e in edges):
                edges.append(x)
    edges.sort()
    if len(edges) < 2:
        return None

    grid = []
    for r in rows:
        cells = [""] * len(edges)
        for c in r:
            idx = min(range(len(edges)),
                      key=lambda i: abs(edges[i] - c["bbox"].x0))
            cells[idx] = (cells[idx] + " " + c["text"]).strip()
        grid.append(cells)

    filled = sum(1 for row in grid for c in row if c)
    if filled < len(edges) * 1.5 or filled / (len(grid) * len(edges)) < 0.35:
        return None
    return grid


def _md_grid_to_md(grid):
    w = len(grid[0])
    def esc(c):
        return c.replace("|", "\\|")
    out = ["| " + " | ".join(esc(c) for c in grid[0]) + " |",
           "|" + "|".join(["---"] * w) + "|"]
    for row in grid[1:]:
        out.append("| " + " | ".join(esc(c) for c in row) + " |")
    return out


def _md_figure_regions(page, table_regions):
    """Figure regions: clustered vector art, plus placed raster images."""
    parea = page.rect.get_area() or 1
    boxes, kinds = [], []

    raw = []
    try:
        for d in page.get_drawings():
            r = fitz.Rect(d["rect"])
            if r.is_empty or r.width < 4 or r.height < 4:
                continue
            if r.get_area() > 0.92 * parea:          # page border / background
                continue
            if any(t.intersects(r) and (t & r).get_area() > 0.6 * r.get_area()
                   for t in table_regions):
                continue
            raw.append(r)
    except Exception:
        raw = []
    for box in _md_cluster_rects(raw, _MD_FIG_GAP):
        members = [r for r in raw if box.intersects(r)]
        if len(members) >= 3 and box.get_area() / parea >= _MD_FIG_MIN_AREA:
            boxes.append(box)
            kinds.append("vector")

    try:
        for im in page.get_images(full=True):
            for r in page.get_image_rects(im[0]):
                if r.get_area() / parea < 0.02:      # logos, rules, bullets
                    continue
                merged = False
                for i, b in enumerate(boxes):
                    if b.intersects(r) and (b & r).get_area() > 0.5 * min(
                            b.get_area(), r.get_area()):
                        boxes[i] = b | r
                        kinds[i] = "mixed"
                        merged = True
                        break
                if not merged:
                    boxes.append(fitz.Rect(r))
                    kinds.append("raster")
    except Exception:
        pass
    return list(zip(boxes, kinds))


def _md_same_column(box, x0, x1):
    """Does this block sit in the column the reader is currently in?"""
    if x1 >= 1e8:                      # end-of-page sweep: anything goes
        return True
    overlap = min(box.x1, x1) - max(box.x0, x0)
    return overlap > min(box.width, max(1.0, x1 - x0)) * 0.35


def _md_bind_caption(box, lines, page_rect, avoid=None):
    """Find the caption for a figure/table region.

    Figure captions usually sit below, table captions above, so both are
    searched — nearest wins, and only text that actually looks like a caption
    is accepted.
    """
    best, best_d = None, 1e9
    for l in lines:
        m = _CAP_RE.match(l["text"])
        if not m:
            continue
        lb = l["bbox"]
        overlap = min(lb.x1, box.x1) - max(lb.x0, box.x0)
        if overlap < min(lb.width, box.width) * 0.25:
            continue
        if lb.y0 >= box.y1:
            d = lb.y0 - box.y1
        elif lb.y1 <= box.y0:
            d = box.y0 - lb.y1
        else:
            d = 0
        if d < best_d and d <= _MD_CAP_GAP:
            best, best_d = l, d
    if not best:
        return None, None, []
    m = _CAP_RE.match(best["text"])
    label = f"{m.group(1).rstrip('.').title()} {m.group(2)}"
    # Caption text may wrap onto following lines.
    parts, taken = [best["text"]], [best]
    for l in lines:
        if l is best:
            continue
        if avoid is not None and avoid.intersects(l["bbox"]) and \
                (avoid & l["bbox"]).get_area() > 0.45 * l["bbox"].get_area():
            continue
        if (abs(l["bbox"].x0 - best["bbox"].x0) < 14
                and 0 <= l["bbox"].y0 - best["bbox"].y1 < best["bbox"].height * 1.9
                and l["size"] <= best["size"] + 0.4
                and not _CAP_RE.match(l["text"])):
            parts.append(l["text"])
            taken.append(l)
            best = l
    return label, _md_tidy(" ".join(parts)), taken


def _md_is_equation(line, body_size):
    t = line["text"]
    if not t or len(t) > 190:
        return False
    if any(ch in t for ch in _MD_MATH):
        letters = sum(ch.isalpha() for ch in t)
        return letters < len(t) * 0.72
    if re.match(r"^\s*\(\d+\)\s*$", t):
        return False
    if line["italic"] and re.search(r"[=<>≤≥]", t) and len(t) < 90:
        return True
    return False


def _md_heading(line, body_size):
    """Heading level, or 0. Size and weight beat pattern matching, but a
    numbered or well-known section name counts even when the size is subtle."""
    t = line["text"]
    if not t or len(t) > 120 or t.endswith((".", ",", ";")) and not _NUM_HEAD_RE.match(t):
        if not _WORD_HEAD_RE.match(t):
            return 0
    big = line["size"] - body_size
    if big >= 5:
        return 1
    if big >= 2.2:
        return 2
    if big >= 0.9 or line["bold"]:
        if _NUM_HEAD_RE.match(t) or _WORD_HEAD_RE.match(t) or line["size"] > body_size:
            depth = t.count(".") if _NUM_HEAD_RE.match(t) else 0
            return min(4, 2 + depth) if depth else 2
    if _WORD_HEAD_RE.match(t) and len(t) < 40:
        return 2
    return 0


def _md_slug(name):
    s = re.sub(r"[^\w\-]+", "_", os.path.splitext(os.path.basename(name))[0])
    return (s.strip("_") or "document")[:48]


def op_pdf2md(raw, name="document.pdf", dpi=170, want_images=True,
              want_tables=True, want_math=True, want_pages=True,
              want_header=True, crop_pad=7, max_figs=120, page_from=None,
              page_to=None):
    """Convert a PDF to Markdown, returning (files, summary_dict).

    files: list of (filename, bytes) — the .md plus any figure crops.
    """
    _md_init()
    try:
        doc = fitz.open(stream=raw, filetype="pdf")
    except Exception as e:                                    # noqa: BLE001
        raise FeatureError(f"Could not open that PDF: {e}")
    if doc.needs_pass:
        raise FeatureError("That PDF is password-protected.")

    stem = _md_slug(name)
    first = max(1, int(page_from or 1))
    last = min(doc.page_count, int(page_to or doc.page_count))
    if first > last:
        raise FeatureError(f"Page range {first}-{last} is empty.")
    idx = list(range(first - 1, last))

    pages_lines, page_rects = [], []
    for i in idx:
        page = doc[i]
        pages_lines.append(_md_lines(_md_spans(page)))
        page_rects.append(page.rect)

    body_size = _md_body_size(pages_lines)
    furniture = _md_furniture(pages_lines, page_rects, body_size)

    out, files, figures, warnings = [], [], [], []
    nfig = 0
    stats = {"pages": len(idx), "figures": 0, "tables": 0, "equations": 0,
             "page_images": 0, "chars": 0, "low_confidence": []}

    for n, (i, lines, rect) in enumerate(zip(idx, pages_lines, page_rects), 1):
        page = doc[i]
        pno = i + 1
        body = [l for l in lines
                if l["size"] > body_size + 1.0
                or re.sub(r"\d+", "#", l["text"])[:70] not in furniture]
        chars = sum(len(l["text"]) for l in body)

        table_regions = _md_table_regions(page) if want_tables else []
        fig_regions = _md_figure_regions(page, table_regions) if want_images else []

        # --- a page with no readable text is a scan: image is the only truth
        if chars < 40:
            if want_pages and want_images:
                pix = page.get_pixmap(dpi=dpi)
                fn = f"{stem}_p{pno:03d}.png"
                files.append((fn, pix.tobytes("png")))
                out += [f"<!-- page {pno}: no text layer; full page image -->",
                        f"![Page {pno} (scanned — no text layer)]({fn})", ""]
                stats["page_images"] += 1
                figures.append({"page": pno, "file": fn, "label": None,
                                "kind": "page-scan", "caption": None})
            else:
                out += [f"<!-- page {pno}: no text layer, image export off -->", ""]
            stats["low_confidence"].append(pno)
            continue

        consumed = set()          # lines already emitted as caption or table

        # --- tables ------------------------------------------------------
        tables = []
        for reg in table_regions:
            grid = _md_table_from(reg, body)
            if not grid:
                continue
            label, cap, cap_lines = _md_bind_caption(reg, body, rect, avoid=reg)
            tables.append({"rect": reg, "grid": grid, "label": label,
                           "caption": cap})
            for l in cap_lines:
                consumed.add(id(l))
            for l in body:
                if reg.intersects(l["bbox"]) and (reg & l["bbox"]).get_area() > \
                        0.45 * l["bbox"].get_area():
                    consumed.add(id(l))

        # --- figures -----------------------------------------------------
        figs = []
        for reg, kind in fig_regions:
            if nfig >= max_figs:
                warnings.append(f"figure cap ({max_figs}) reached at page {pno}")
                break
            if any(t["rect"].intersects(reg) and
                   (t["rect"] & reg).get_area() > 0.6 * reg.get_area()
                   for t in tables):
                continue
            label, cap, cap_lines = _md_bind_caption(reg, body, rect, avoid=reg)
            for l in cap_lines:
                consumed.add(id(l))
            clip = fitz.Rect(reg.x0 - crop_pad, reg.y0 - crop_pad,
                             reg.x1 + crop_pad, reg.y1 + crop_pad) & page.rect
            try:
                pix = page.get_pixmap(dpi=dpi, clip=clip)
                data = pix.tobytes("png")
            except Exception:                                 # noqa: BLE001
                continue
            nfig += 1
            fn = f"{stem}_p{pno:03d}_fig{nfig:02d}.png"
            files.append((fn, data))
            figs.append({"rect": reg, "file": fn, "label": label,
                         "caption": cap, "kind": kind})
            figures.append({"page": pno, "file": fn, "label": label,
                            "kind": kind, "caption": cap})
        stats["figures"] += len(figs)
        stats["tables"] += len(tables)

        # --- flow the page in reading order ------------------------------
        cuts = _md_columns(body, rect)
        ordered = _md_order(body, rect, cuts)
        out.append(f"<!-- page {pno} -->")

        placed_t, placed_f = set(), set()
        para, last_line = [], None

        def flush():
            nonlocal out
            if para:
                out.append(" ".join(para))
                out.append("")
                para.clear()

        def emit_blocks_before(y, x0, x1):
            """Drop in any figure/table whose top sits above this line."""
            nonlocal out
            for ti, t in enumerate(tables):
                if ti in placed_t or t["rect"].y0 > y:
                    continue
                if not _md_same_column(t["rect"], x0, x1):
                    continue
                flush()
                cap = t["caption"] or (f"{t['label']}" if t["label"] else None)
                if cap:
                    out += [f"**{cap}**", ""]
                out += _md_grid_to_md(t["grid"]) + [""]
                placed_t.add(ti)
            for fi, f in enumerate(figs):
                if fi in placed_f or f["rect"].y0 > y:
                    continue
                if not _md_same_column(f["rect"], x0, x1):
                    continue
                flush()
                alt = f["label"] or f"Figure on page {pno}"
                out.append(f"![{alt}]({f['file']})")
                if f["caption"]:
                    out += ["", f"*{f['caption']}*"]
                elif not f["label"]:
                    out += ["", f"*(unlabelled {f['kind']} figure, page {pno})*"]
                out.append("")
                placed_f.add(fi)

        for l in ordered:
            if id(l) in consumed:
                continue
            emit_blocks_before(l["bbox"].y0, l["bbox"].x0, l["bbox"].x1)
            t = l["text"]
            if not t:
                continue
            lvl = _md_heading(l, body_size)
            if lvl:
                flush()
                out += ["#" * lvl + " " + t, ""]
                last_line = l
                continue
            if want_math and _md_is_equation(l, body_size):
                flush()
                out += ["$$", t, "$$", ""]
                stats["equations"] += 1
                last_line = l
                continue
            if l["size"] < body_size - 1.0 and _FOOT_RE.match(t):
                flush()
                note = [t]
                for nxt in ordered[ordered.index(l) + 1:]:
                    if (id(nxt) in consumed or nxt["size"] > body_size - 1.0
                            or abs(nxt["bbox"].x0 - l["bbox"].x0) > 30
                            or nxt["bbox"].y0 - l["bbox"].y1 > l["bbox"].height * 1.6
                            or _FOOT_RE.match(nxt["text"])):
                        break
                    note.append(nxt["text"])
                    consumed.add(id(nxt))
                    l = nxt
                out += ["> " + " ".join(note), ""]
                last_line = l
                continue
            # Decide the break BEFORE appending: a wide vertical gap, or a jump
            # to another column, ends the previous paragraph rather than this one.
            if last_line is not None and para:
                gap = l["bbox"].y0 - last_line["bbox"].y1
                jumped = (l["bbox"].y0 < last_line["bbox"].y0 - 2
                          or abs(l["bbox"].x0 - last_line["bbox"].x0) > 40)
                if gap > l["bbox"].height * 1.4 or jumped:
                    flush()
            # join wrapped lines; a hyphen at end of line means a split word
            if para and para[-1].endswith("-") and not para[-1].endswith("--"):
                para[-1] = para[-1][:-1] + t.lstrip()
            else:
                para.append(t)
            last_line = l
        flush()
        emit_blocks_before(1e9, 0, 1e9)          # anything left at page end
        out.append("")

        # --- confidence --------------------------------------------------
        if want_images and chars < 220 and (figs or fig_regions):
            stats["low_confidence"].append(pno)
        stats["chars"] += chars

    doc.close()

    # ------------------------------------------------------------------ #
    #  Escalate genuinely unreliable pages to a full page image.
    # ------------------------------------------------------------------ #
    body_md = "\n".join(out)
    body_md = re.sub(r"\n{3,}", "\n\n", body_md).strip() + "\n"

    head = []
    if want_header:
        head = _md_header(name, stem, stats, figures, len(idx), first, last)
    md = ("\n".join(head) + "\n" + body_md) if head else body_md
    files.insert(0, (f"{stem}.md", md.encode("utf-8")))

    summary = dict(stats)
    summary.update({
        "stem": stem, "md_name": f"{stem}.md", "images": len(files) - 1,
        "md_bytes": len(md.encode()), "figure_list": figures,
        "warnings": warnings,
    })
    return files, summary


def _md_header(name, stem, stats, figures, npages, first, last):
    """A short self-describing preamble, so the file needs no covering note."""
    rng = "" if (first == 1 and last == npages + first - 1) else \
        f", pages {first}-{last}"
    h = [f"# {os.path.splitext(os.path.basename(name))[0]}", "",
         "> **How to read this file.** It was converted from a PDF"
         f" (`{os.path.basename(name)}`{rng}).",
         "> Text is real extracted text. Images are cropped from the page and"
         " appear at the point",
         "> in the reading order where they occur, so the figure above a caption"
         " is the figure that",
         "> caption describes. `<!-- page N -->` marks where each PDF page"
         " begins."]
    if stats["figures"] or stats["page_images"]:
        h.append(f"> Figures are separate PNG files next to this one"
                 f" (`{stem}_pNNN_figNN.png`); keep them together.")
    if stats["low_confidence"]:
        pages = ", ".join(str(p) for p in stats["low_confidence"][:14])
        more = "…" if len(stats["low_confidence"]) > 14 else ""
        h.append(f"> **Lower confidence on page(s) {pages}{more}** — little or no"
                 " extractable text there, so trust the image over the text.")
    h.append("")
    return h


def op_pdf2md_job(got, data):
    """Adapter: run a conversion and shape it for the download endpoint."""
    files, summary = op_pdf2md(
        got["bytes"], got["name"],
        dpi=int(data.get("dpi", 170)),
        want_images=bool(data.get("images", True)),
        want_tables=bool(data.get("tables", True)),
        want_math=bool(data.get("math", True)),
        want_header=bool(data.get("header", True)),
        page_from=data.get("page_from"), page_to=data.get("page_to"))
    if data.get("index"):
        entry = {"source": got["name"],
                 **{k: summary[k] for k in ("md_name", "pages", "figures",
                                            "tables", "page_images",
                                            "low_confidence", "images")}}
        files.append(("_INDEX.md", md_index([entry]).encode()))
    bits = [f"{summary['pages']} page(s)", f"{summary['figures']} figure(s)",
            f"{summary['tables']} table(s)"]
    if summary["page_images"]:
        bits.append(f"{summary['page_images']} page image(s)")
    note = " · ".join(bits) + f" · {summary['md_bytes']/1024:.0f} KB markdown"
    if summary["low_confidence"]:
        note += (" · check page(s) "
                 + ", ".join(map(str, summary["low_confidence"][:8])))
    return files, note


def md_index(entries, title="PDF → Markdown conversion index"):
    """Optional index for a batch. This is for the human auditing a run — a
    single converted paper does not need it, because its own header explains
    itself."""
    now = time.strftime("%Y-%m-%d %H:%M")
    out = [f"# {title}", "", f"Generated {now}", "",
           "| Source | Markdown | Pages | Figures | Tables | Needs eyes |",
           "|---|---|---|---|---|---|"]
    weak = []
    for e in entries:
        if e.get("error"):
            out.append(f"| {e['source']} | **failed** | - | - | - | {e['error']} |")
            continue
        lc = e.get("low_confidence") or []
        out.append(
            f"| {e['source']} | `{e['md_name']}` | {e['pages']} | "
            f"{e['figures']} | {e['tables']} | "
            f"{'p' + ', p'.join(map(str, lc)) if lc else '—'} |")
        if lc:
            weak.append((e["source"], lc))
    out += ["", "## Pages worth opening yourself", ""]
    if weak:
        for src, lc in weak:
            out.append(f"- **{src}** — page(s) {', '.join(map(str, lc))}: little or "
                       "no extractable text, so the image is the only record.")
    else:
        out.append("- (none — every page had a usable text layer)")
    out += ["", "## Reading these files", "",
            "- Each `.md` opens with a short note explaining its own layout, so it",
            "  can be handed to an LLM on its own with no covering message.",
            "- Figures are cropped PNGs beside each `.md`. Keep them in the same",
            "  folder or the image links break.",
            "- An image sits at the point in the text where it appears in the PDF,",
            "  and its caption follows it, so figure and caption cannot be mixed up.",
            ""]
    return "\n".join(out)


def md_batch(paths, outdir, index=True, **kw):
    """Convert many PDFs to <outdir>/<stem>/. Returns the list of entries."""
    entries = []
    os.makedirs(outdir, exist_ok=True)
    for path in paths:
        src = os.path.basename(path)
        try:
            with open(path, "rb") as fh:
                raw = fh.read()
            files, summary = op_pdf2md(raw, src, **kw)
        except Exception as e:                                # noqa: BLE001
            entries.append({"source": src, "error": str(e)[:120]})
            print(f"  {src:40s} FAILED: {str(e)[:60]}")
            continue
        sub = os.path.join(outdir, summary["stem"])
        os.makedirs(sub, exist_ok=True)
        for fn, data in files:
            with open(os.path.join(sub, fn), "wb") as fh:
                fh.write(data)
        entry = {"source": src, **{k: summary[k] for k in
                 ("md_name", "pages", "figures", "tables", "page_images",
                  "low_confidence", "images")}}
        entries.append(entry)
        print(f"  {src:40s} {summary['pages']:3d}p  {summary['figures']:2d} fig  "
              f"{summary['tables']:2d} tbl  -> {summary['stem']}/"
              f"{'  (check pages ' + ','.join(map(str, summary['low_confidence'])) + ')' if summary['low_confidence'] else ''}")
    if index:
        with open(os.path.join(outdir, "_INDEX.md"), "w") as fh:
            fh.write(md_index(entries))
    return entries


def md_cli(argv):
    """`python multi_toolkit.py md <pdf-or-folder> [outdir] [options]`"""
    import argparse
    ap = argparse.ArgumentParser(
        prog="multi_toolkit.py md",
        description="Convert PDFs to Markdown with cropped figures, for feeding "
                    "to an LLM without shipping the whole PDF.")
    ap.add_argument("src", help="a PDF file, or a folder of PDFs")
    ap.add_argument("outdir", nargs="?", default="./md_out")
    ap.add_argument("--dpi", type=int, default=170, help="figure raster dpi")
    ap.add_argument("--no-images", action="store_true", help="text only")
    ap.add_argument("--no-tables", action="store_true")
    ap.add_argument("--no-math", action="store_true")
    ap.add_argument("--no-header", action="store_true",
                    help="omit the self-describing preamble")
    ap.add_argument("--index", action="store_true",
                    help="also write _INDEX.md summarising the run")
    ap.add_argument("--pages", metavar="A-B", help="page range, e.g. 3-12")
    a = ap.parse_args(argv)

    src = os.path.expanduser(a.src)
    if os.path.isdir(src):
        paths = sorted(os.path.join(r, f) for r, _d, fs in os.walk(src)
                       for f in fs if f.lower().endswith(".pdf"))
        if not paths:
            sys.exit(f"No PDFs found in {src}")
    elif os.path.isfile(src):
        paths = [src]
    else:
        sys.exit(f"No such file or folder: {src}")

    pf = pt = None
    if a.pages:
        m = re.match(r"^\s*(\d+)\s*[-:]\s*(\d+)\s*$", a.pages)
        if not m:
            sys.exit("--pages wants a range like 3-12")
        pf, pt = int(m.group(1)), int(m.group(2))

    print(f"Converting {len(paths)} PDF(s) -> {a.outdir}")
    md_batch(paths, os.path.expanduser(a.outdir), index=a.index, dpi=a.dpi,
             want_images=not a.no_images, want_tables=not a.no_tables,
             want_math=not a.no_math, want_header=not a.no_header,
             page_from=pf, page_to=pt)
    print(f"\nDone -> {a.outdir}")
    if a.index:
        print(f"Index: {os.path.join(a.outdir, '_INDEX.md')}")


# --------------------------------------------------------------------------- #
#  HTTP server
# --------------------------------------------------------------------------- #
class Handler(BaseHTTPRequestHandler):
    protocol_version = "HTTP/1.1"  # honor Content-Length for large bodies

    def _raw(self, code, body, ctype="text/plain; charset=utf-8", extra=None):
        if isinstance(body, str):
            body = body.encode("utf-8")
        self.send_response(code)
        self.send_header("Content-Type", ctype)
        self.send_header("Content-Length", str(len(body)))
        for k, v in (extra or {}).items():
            # HTTP headers must be Latin-1; never let an exotic char crash the body write.
            safe = str(v).encode("latin-1", "replace").decode("latin-1")
            self.send_header(k, safe)
        self.end_headers()
        self.wfile.write(body)

    def _json(self, obj, code=200):
        self._raw(code, json.dumps(obj), "application/json")

    def _deliver(self, results, info):
        """Send a single file, or zip multiple results."""
        from urllib.parse import quote
        if not results:
            self._raw(400, "Nothing was produced.")
            return
        if len(results) == 1:
            name, data = results[0]
            ascii_name = name.encode("ascii", "ignore").decode() or "download"
            self._raw(200, data, _guess_ctype(name),
                      extra={"X-Filename": quote(name), "X-Info": quote(info),
                             "Content-Disposition":
                                 f"attachment; filename=\"{ascii_name}\"; "
                                 f"filename*=UTF-8''{quote(name)}"})
        else:
            buf = io.BytesIO()
            with zipfile.ZipFile(buf, "w", zipfile.ZIP_DEFLATED) as z:
                for name, data in results:
                    z.writestr(name, data)
            self._raw(200, buf.getvalue(), "application/zip",
                      extra={"X-Filename": "results.zip",
                             "X-Info": quote(info + f"\n({len(results)} files, zipped)"),
                             "Content-Disposition": 'attachment; filename="results.zip"'})

    def do_GET(self):
        if self.path in ("/", "/index.html"):
            self._raw(200, PAGE, "text/html; charset=utf-8")
        elif self.path == "/capabilities":
            self._json({"libreoffice": bool(office_binary()),
                        "ghostscript": bool(gs_binary()),
                        "ffmpeg": bool(ffmpeg_binary())})
        elif self.path.startswith("/yt_progress"):
            job = _YT_JOBS.get(self._query().get("id", ""))
            if not job:
                self._raw(404, "Unknown download job.")
                return
            with _YT_LOCK:
                pub = {k: job.get(k) for k in
                       ("status", "pct", "msg", "speed", "eta", "name", "detail", "rate", "error")}
            self._json(pub)
        elif self.path.startswith("/yt_file"):
            self._yt_send(self._query().get("id", ""))
        else:
            self._raw(404, "Not found")

    def _query(self):
        from urllib.parse import urlparse, parse_qs
        q = parse_qs(urlparse(self.path).query)
        return {k: v[0] for k, v in q.items()}

    def _yt_send(self, job_id):
        """Stream a finished video from disk (never buffered in RAM)."""
        from urllib.parse import quote
        job = _YT_JOBS.get(job_id)
        if not job or job.get("status") != "done" or not job.get("path"):
            self._raw(404, "No finished download with that id.")
            return
        path, name = job["path"], job["name"]
        try:
            size = os.path.getsize(path)
            ascii_name = name.encode("ascii", "ignore").decode() or "video.mp4"
            self.send_response(200)
            self.send_header("Content-Type", _guess_ctype(name))
            self.send_header("Content-Length", str(size))
            self.send_header("Content-Disposition",
                             f"attachment; filename=\"{ascii_name}\"; "
                             f"filename*=UTF-8''{quote(name)}")
            self.end_headers()
            with open(path, "rb") as fh:
                shutil.copyfileobj(fh, self.wfile, 256 * 1024)
        finally:
            _yt_cleanup(job_id)

    def do_POST(self):
        try:
            length = int(self.headers.get("Content-Length", 0))
            data = json.loads(self.rfile.read(length)) if length else {}

            if self.path == "/inspect":
                files = resolve_files([data.get("file")])
                self._json(op_inspect(files[0], data.get("thumbs", "cover")))
                return

            if self.path == "/yt_info":
                self._json(op_yt_info(data.get("url", "").strip(),
                                      data.get("cookies", "")))
                return

            if self.path == "/yt_start":
                self._json(op_yt_start(data.get("url", ""), data.get("height"),
                                       data.get("compat", "best"),
                                       data.get("vertical", "off"),
                                       data.get("vsize", "1080"),
                                       data.get("cookies", "")))
                return

            if self.path == "/qr":
                payload = data.get("data")
                if payload is None:
                    payload = qr_payload(data.get("kind", "text"),
                                         data.get("fields", {}))
                self._json(op_qr(
                    data=payload, ec=data.get("ec", "H"),
                    target_px=data.get("target_px", 1024),
                    border=data.get("border", 4), fg=data.get("fg", "#000000"),
                    bg=data.get("bg", "#ffffff"), style=data.get("style", "square"),
                    logo_data=data.get("logo_data"),
                    logo_pct=data.get("logo_pct", 22),
                    logo_style=data.get("logo_style", "original"),
                    pad=bool(data.get("pad", True)),
                    pad_shape=data.get("pad_shape", "rounded"),
                    pad_pct=data.get("pad_pct", 6),
                    fmt=data.get("fmt", "png")))
                return

            if self.path == "/img_fetch":
                if data.get("url"):
                    raw, name = op_img_fetch(data["url"])
                else:
                    got = resolve_files([data.get("file")])
                    raw, name = got[0]["bytes"], got[0]["name"]
                token = cache_put(name, raw)
                out = img_describe(raw, name)
                out["token"] = token
                self._json(out)
                return

            if self.path == "/merge":
                files = resolve_files(data.get("files", []))
                if not files:
                    raise FeatureError("Add at least one PDF to merge.")
                results, info = op_merge(files, bool(data.get("bookmarks", True)))
            elif self.path == "/split":
                f = data.get("file")
                if not f:
                    raise FeatureError("Add a PDF to split.")
                files = resolve_files([f])
                results, info = op_split(files[0], data.get("mode", "each"),
                                         data.get("ranges", ""),
                                         data.get("every_n", 1),
                                         data.get("pages"))
            elif self.path == "/compress":
                files = resolve_files(data.get("files", []))
                if not files:
                    raise FeatureError("Add at least one PDF to compress.")
                p = COMPRESS_PRESETS.get(data.get("preset"))
                if p:
                    results, info = op_compress(files, p["mode"], p["quality"], p["max_px"])
                else:
                    results, info = op_compress(files, data.get("mode", "strong"),
                                                data.get("quality", 60), data.get("max_px", 1600))
            elif self.path == "/compress_preview":
                files = resolve_files(data.get("files", []))
                if not files:
                    raise FeatureError("Add at least one PDF to compress.")
                self._json(op_compress_preview(files))
                return
            elif self.path == "/download":
                files = resolve_files(data.get("files", []))
                results = [(f["name"], f["bytes"]) for f in files]
                info = data.get("info", "Done.")
            elif self.path == "/pdf2md":
                got = resolve_files([data.get("file")])
                results, info = op_pdf2md_job(got[0], data)
            elif self.path == "/img_process":
                got = resolve_files([data.get("file")])
                results, info = op_img_process(
                    got[0]["bytes"], got[0]["name"], crop=data.get("crop"),
                    out_w=data.get("out_w"), out_h=data.get("out_h"),
                    fmt=data.get("fmt", "png"), quality=data.get("quality", 90),
                    rotate=data.get("rotate", 0), flip=bool(data.get("flip")),
                    grayscale=bool(data.get("grayscale")))
            elif self.path == "/convert":
                route = data.get("route")
                fn = CONVERT_ROUTES.get(route)
                if not fn:
                    raise FeatureError(f"Unknown conversion: {route}")
                files = resolve_files(data.get("files", []))
                if not files:
                    raise FeatureError("Add at least one file.")
                results, info = fn(files, img_format=data.get("img_format", "png"),
                                   dpi=data.get("dpi", 150))
            else:
                self._raw(404, "Not found")
                return

            self._deliver(results, info)

        except CacheMiss as cm:
            self._raw(409, f"cache-miss:{cm}")
        except FeatureError as fe:
            self._raw(400, str(fe))
        except Exception:  # noqa: BLE001
            self._raw(500, "Server error:\n" + traceback.format_exc())

    def log_message(self, *args):
        pass


def _guess_ctype(name):
    ext = os.path.splitext(name)[1].lower()
    return {
        ".pdf": "application/pdf",
        ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        ".png": "image/png", ".jpg": "image/jpeg", ".jpeg": "image/jpeg",
        ".zip": "application/zip",
        ".mp4": "video/mp4", ".webm": "video/webm", ".mkv": "video/x-matroska",
        ".m4a": "audio/mp4",
    }.get(ext, "application/octet-stream")


def _free_port(preferred=8000):
    for port in (preferred, 8001, 8080, 0):
        try:
            s = socket.socket()
            s.bind(("127.0.0.1", port))
            p = s.getsockname()[1]
            s.close()
            return p
        except OSError:
            continue
    return 0


# --------------------------------------------------------------------------- #
#  Front-end
# --------------------------------------------------------------------------- #
PAGE = r"""<!doctype html>
<html lang="en"><head><meta charset="utf-8">
<meta name="viewport" content="width=device-width, initial-scale=1">
<title>Multi Toolkit</title>
<style>
/* ==========================================================================
   Multi Toolkit — "bench instrument"
   Chassis in warm graphite. Structure is silkscreened in mono; content is set
   in the UI face. Two accents with strict jobs: AMBER is the machine reporting
   a measured fact, STEEL is something you can act on. Nothing else is coloured.
   System fonts only — this runs offline, so no webfont may be required.
   ========================================================================== */
:root{
  --void:#100E0C;      /* chassis ground — warm black, never blue-black       */
  --bench:#191512;     /* panel surface                                       */
  --riser:#221D19;     /* raised control                                      */
  --edge:#342C26;      /* machined hairline                                   */
  --edge-hi:#4A3F36;
  --paper:#EFE9E1;     /* content                                             */
  --dim:#9C8F84;       /* secondary                                           */
  --faint:#827568;     /* tertiary / disabled — 4.06:1 on bench, kept legible */
  --amber:#FFB000;     /* LIVE measured data + progress. Nothing else.        */
  --amber-dim:#8A6410;
  --steel:#7FAECF;     /* interactive / selected. Nothing else.               */
  --steel-dim:#2C4557;
  --rust:#DC6148;      /* destructive + error. Nothing else. 5.05:1 on bench  */
  --ok:#8FBF7A;
  --amber-lit:#FFB000; /* fixed: only inside the always-dark readout           */
  --steel-ink:#CFE6F5; /* text on a steel-tinted fill                          */
  --on-steel:#0B1219;  /* text on a solid steel fill                           */
  --on-amber:#100E0C;  /* text on a solid amber fill                           */
  --ok-edge:#31402C; --err-edge:#4A2A24; --err-bg:#2A1A16; --err-ink:#E9A192;
  --warn-bg:#241B0E; --warn-edge:#4A3A16; --warn-ink:#F0C77E;
  --veil:#141110F2;

  --ui:-apple-system,BlinkMacSystemFont,"Segoe UI",Roboto,Helvetica,Arial,sans-serif;
  --mono:ui-monospace,"SF Mono","JetBrains Mono","Cascadia Mono",Menlo,Consolas,monospace;

  --rail-w:214px;
  --r:4px;             /* machined, not pillowy */
  --r-lg:7px;
}
*{box-sizing:border-box}
/* Author styles beat the UA sheet, so a rule like .row{display:flex} would
   otherwise override the built-in [hidden]{display:none}. Restate it. */
[hidden]{display:none!important}
::selection{background:var(--steel);color:var(--void)}
html{color-scheme:dark}
body{
  margin:0;background:var(--void);color:var(--paper);
  font:15px/1.55 var(--ui);
  font-variant-numeric:tabular-nums;
  -webkit-font-smoothing:antialiased;
}
/* Every number in this interface lines up in a column. */
.mono,.readout,.pill,.tab,.chip,.eyebrow,.status,.qrstat,.ytnote,.hint,
.pv-sum,.sz,.ro-cell i,.fld>span,.card h4,.group span,.d{
  font-variant-numeric:tabular-nums}

:focus-visible{outline:2px solid var(--steel);outline-offset:2px;border-radius:2px}
.hide{display:none!important}

/* ---------- masthead ---------------------------------------------------- */
.masthead{
  display:flex;align-items:center;gap:18px;flex-wrap:wrap;
  padding:13px 20px;background:var(--bench);
  border-bottom:1px solid var(--edge);
  position:sticky;top:0;z-index:40;
}
.brand{display:flex;align-items:baseline;gap:11px;margin-right:auto}
.brand .mark{
  font:700 15px/1 var(--mono);color:var(--on-amber);background:var(--amber);
  padding:5px 6px;border-radius:2px;letter-spacing:-.06em;align-self:center}
.brand h1{
  margin:0;font-size:16px;font-weight:650;letter-spacing:-.015em;white-space:nowrap}
.eyebrow{
  font:600 10px/1 var(--mono);letter-spacing:.2em;text-transform:uppercase;
  color:var(--faint)}
.pills{display:flex;gap:6px;flex-wrap:wrap}
.pill{
  font:600 10px/1 var(--mono);letter-spacing:.11em;text-transform:uppercase;
  padding:6px 9px;border-radius:2px;border:1px solid var(--edge);
  color:var(--faint);background:var(--void);white-space:nowrap;
  display:inline-flex;align-items:center;gap:6px}
.pill::before{
  content:"";width:5px;height:5px;border-radius:50%;background:currentColor;
  opacity:.55}
.pill.on{color:var(--ok);border-color:var(--ok-edge)}
.pill.off{color:var(--faint)}

/* ---------- shell ------------------------------------------------------- */
.shell{display:grid;grid-template-columns:var(--rail-w) minmax(0,1fr);
  align-items:start}
.rail{
  position:sticky;top:53px;align-self:start;
  height:calc(100vh - 53px);overflow-y:auto;
  padding:16px 10px 22px;border-right:1px solid var(--edge);
  background:linear-gradient(180deg,var(--bench),var(--void) 220px)}
.stage{min-width:0;padding:20px 22px 132px;max-width:1180px}

/* ---------- rail navigation --------------------------------------------- */
.groups{margin:0 0 7px;padding:0 8px}
.groups + .tabs{margin-bottom:20px}
.group{
  display:flex;align-items:baseline;gap:7px;cursor:pointer;user-select:none;
  padding:6px 0 5px;border-bottom:1px solid var(--edge)}
.group b{
  font:600 10px/1 var(--mono);letter-spacing:.19em;text-transform:uppercase;
  color:var(--dim);transition:color .12s}
.group span{
  font:10px/1 var(--mono);letter-spacing:.08em;color:var(--faint);
  margin-left:auto;opacity:.75}
.group:hover b{color:var(--paper)}
.group.active b{color:var(--amber)}
.tabs{display:flex;flex-direction:column;gap:1px}
.tab{
  position:relative;cursor:pointer;user-select:none;
  padding:8px 11px 8px 14px;border-radius:var(--r);
  font-size:14px;color:var(--dim);
  transition:background .12s,color .12s}
.tab::before{
  content:"";position:absolute;left:4px;top:50%;transform:translateY(-50%);
  width:2px;height:0;background:var(--steel);border-radius:1px;
  transition:height .14s ease}
.tab:hover{background:var(--riser);color:var(--paper)}
.tab.active{background:var(--riser);color:var(--paper);font-weight:600}
.tab.active::before{height:15px}

/* ---------- THE READOUT — the one loud thing ---------------------------- */
.readout{
  display:flex;flex-wrap:wrap;gap:0 30px;align-items:flex-end;
  padding:13px 17px;margin-bottom:18px;
  background:
    linear-gradient(180deg,#0C0A08,#080706);
  border:1px solid var(--edge);border-radius:var(--r-lg);
  box-shadow:inset 0 1px 0 #ffffff0a, inset 0 0 34px #FFB0000a;
}
.ro-cell{display:flex;flex-direction:column-reverse;gap:3px;min-width:0}
.ro-cell i{
  font:600 9.5px/1 var(--mono);letter-spacing:.19em;text-transform:uppercase;
  color:var(--faint);font-style:normal}
.ro-cell b{
  font:500 19px/1.05 var(--mono);letter-spacing:-.01em;color:var(--amber-lit);
  text-shadow:0 0 16px #FFB00040;white-space:nowrap}
.readout[data-state="idle"] .ro-cell b{
  color:var(--faint);text-shadow:none;font-size:15px;font-weight:400}

/* ---------- surfaces ---------------------------------------------------- */
.card,.ytcard,.ytbox,.preview,.presets{
  background:var(--bench);border:1px solid var(--edge);border-radius:var(--r-lg)}
.card{padding:15px 16px;margin-bottom:12px}
.card h4{
  margin:0 0 11px;font:600 10px/1 var(--mono);letter-spacing:.19em;
  text-transform:uppercase;color:var(--faint)}

/* ---------- drop zone ---------------------------------------------------- */
.drop{
  border:1px dashed var(--edge-hi);border-radius:var(--r-lg);
  background:var(--bench);cursor:pointer;text-align:center;
  display:flex;flex-direction:column;align-items:center;justify-content:center;
  gap:7px;transition:border-color .13s,background .13s,color .13s}
.drop.big{padding:46px 24px}
.drop.slim{padding:15px 20px;flex-direction:row;gap:12px}
.drop.slim .hint{display:none}
.drop b{color:var(--paper);font-weight:600}
.drop{color:var(--dim)}
.drop:hover,.drop.hot{border-color:var(--steel);background:var(--riser);color:var(--paper)}
.drop .hint{font-size:12.5px;color:var(--faint);max-width:46ch;line-height:1.5}

/* ---------- file list ---------------------------------------------------- */
.files{list-style:none;margin:12px 0 0;padding:0;display:flex;
  flex-direction:column;gap:1px}
.files li{
  display:flex;align-items:center;gap:12px;padding:9px 12px;
  background:var(--bench);border:1px solid var(--edge);border-radius:var(--r);
  transition:border-color .12s,background .12s}
.files li:hover{border-color:var(--edge-hi)}
.files li.dragging{opacity:.45;border-color:var(--steel)}
.files li.reorderable{cursor:grab}
.files li.reorderable:active{cursor:grabbing}
.files .pv-title,.files .t{font-size:14px;min-width:0;overflow:hidden;
  text-overflow:ellipsis;white-space:nowrap}
.files .d,.sz{font:11.5px var(--mono);color:var(--faint);white-space:nowrap}
.files img,.files canvas{
  width:34px;height:44px;object-fit:cover;border-radius:2px;
  background:var(--riser);border:1px solid var(--edge);flex:none}

/* ---------- controls ----------------------------------------------------- */
.btn,button.primary,.mini{
  font:inherit;cursor:pointer;border-radius:var(--r);transition:.12s;
  border:1px solid var(--edge-hi);background:var(--riser);color:var(--paper)}
.btn{padding:8px 14px;font-size:13.5px}
.btn:hover{background:var(--edge);border-color:var(--edge-hi)}
.btn.danger{color:var(--rust);border-color:var(--err-edge)}
.btn.danger:hover{background:var(--err-bg);border-color:var(--rust)}
button.primary{
  padding:12px 26px;font-size:14.5px;font-weight:650;letter-spacing:.01em;
  background:var(--steel);color:var(--on-steel);border-color:var(--steel)}
button.primary:hover:not(:disabled){background:#98C1DE;border-color:#98C1DE}
button.primary:disabled{
  background:var(--riser);color:var(--faint);border-color:var(--edge);
  cursor:not-allowed}
input[type=text],input[type=number],input[type=url],select,textarea{
  background:var(--void);color:var(--paper);
  border:1px solid var(--edge-hi);border-radius:var(--r);
  padding:8px 10px;font:inherit;font-size:14px}
input:focus,select:focus,textarea:focus{outline:none;border-color:var(--steel)}
input:disabled,select:disabled{opacity:.45;cursor:not-allowed}
input[type=checkbox]{accent-color:var(--steel);width:15px;height:15px}
input[type=range]{accent-color:var(--amber);width:118px}
input[type=color]{
  width:42px;height:30px;padding:0;border:1px solid var(--edge-hi);
  border-radius:var(--r);background:none;cursor:pointer}
.swatch{width:42px;height:30px;padding:0;border:1px solid var(--edge-hi);
  border-radius:var(--r);background:none;cursor:pointer}

/* segmented chips — a chosen option, so STEEL */
.chips{display:flex;gap:5px;flex-wrap:wrap;margin-bottom:12px}
.chip{
  font:600 11.5px/1 var(--mono);letter-spacing:.06em;
  padding:7px 11px;border-radius:var(--r);cursor:pointer;user-select:none;
  border:1px solid var(--edge);background:var(--bench);color:var(--dim);
  transition:.12s}
.chip:hover{color:var(--paper);border-color:var(--edge-hi)}
.chip.active{
  background:var(--steel-dim);border-color:var(--steel);color:var(--steel-ink)}

.opts{
  display:flex;gap:16px;flex-wrap:wrap;align-items:center;
  padding:13px 15px;margin-bottom:12px;
  background:var(--bench);border:1px solid var(--edge);border-radius:var(--r-lg)}
.opts label,.row label{
  display:flex;align-items:center;gap:8px;font-size:13.5px;color:var(--dim)}
.opts select,.opts input[type=text],.opts input[type=number]{color:var(--paper)}
.row{display:flex;gap:14px;flex-wrap:wrap;align-items:center}
.rng{display:flex;align-items:center;gap:9px}
.rng b{font:600 12px var(--mono);color:var(--amber);min-width:40px}
.toolbar{display:flex;gap:8px;align-items:center;margin-top:12px}
.sp{flex:1}

/* ---------- split page grid ---------------------------------------------- */
.preview{padding:14px;margin-bottom:12px}
.grid{display:flex;flex-wrap:wrap;gap:0;align-items:flex-start}
.pg{
  position:relative;width:82px;cursor:pointer;border-radius:var(--r);
  padding:5px;transition:background .12s}
.pg:hover{background:var(--riser)}
.pg img,.pg canvas{
  width:100%;display:block;border-radius:2px;border:1px solid var(--edge);
  background:#fff}
.pg .n{
  display:block;text-align:center;font:11px var(--mono);color:var(--faint);
  margin-top:4px}
.pg.active img,.pg.active canvas{border-color:var(--steel);
  box-shadow:0 0 0 1px var(--steel)}
.pg.wait{opacity:.4}
.cut{
  width:16px;align-self:stretch;cursor:pointer;position:relative;
  border-radius:2px;transition:background .12s}
.cut::after{
  content:"";position:absolute;left:50%;top:8px;bottom:26px;width:1px;
  transform:translateX(-50%);background:var(--edge-hi);transition:.12s}
.cut:hover::after{background:var(--steel);width:2px}
.cut.on::after{background:var(--amber);width:2px;box-shadow:0 0 8px #FFB00066}
.pv-head{display:flex;align-items:center;gap:12px;margin-bottom:11px}
.pv-title{font-weight:600;font-size:14px}
.pv-sum,.pv-note{font:11.5px var(--mono);color:var(--faint)}
.pv-sum{margin-left:auto}
.pv-note{margin-top:9px;line-height:1.5}

/* ---------- compress presets --------------------------------------------- */
.presets{display:grid;grid-template-columns:repeat(auto-fit,minmax(150px,1fr));
  gap:1px;padding:1px;margin-bottom:12px;background:var(--edge);
  border-radius:var(--r-lg);overflow:hidden;border:1px solid var(--edge)}
.pcard{
  background:var(--bench);padding:13px 14px;cursor:pointer;transition:.12s}
.pcard:hover{background:var(--riser)}
.pcard.active{background:var(--steel-dim);box-shadow:inset 2px 0 0 var(--steel)}
.pcard .t{display:block;font-weight:600;font-size:13.5px;margin-bottom:3px}
.pcard .d{display:block;font-size:11.5px;line-height:1.45;color:var(--faint)}
.pcard .sz{display:block;margin-top:7px;font:600 12px var(--mono);color:var(--amber)}
.pcard.active .sz{color:var(--amber)}

/* ---------- video panel --------------------------------------------------- */
.ytbox{padding:15px 16px;margin-bottom:12px}
.ytrow{display:flex;gap:9px;margin-bottom:13px}
.ytrow input{flex:1;min-width:0}
.ytcard{display:flex;gap:14px;padding:13px;margin-bottom:13px;align-items:flex-start}
.ytcard img{width:150px;border-radius:var(--r);border:1px solid var(--edge);flex:none}
.yttitle{font-weight:650;font-size:15px;line-height:1.35;margin-bottom:5px}
.ytsub,.ytmeta{font:11.5px var(--mono);color:var(--faint);letter-spacing:.04em}
.ytmeta{margin-top:4px}
.ytsel{display:flex;gap:14px;flex-wrap:wrap;align-items:center;margin-top:11px}
.ytnote{font-size:12.5px;line-height:1.55;color:var(--faint);margin-top:4px}

/* ---------- QR panel ------------------------------------------------------ */
.qrgrid{display:grid;grid-template-columns:minmax(0,1fr) 296px;gap:14px;
  align-items:start}
.qrfields{display:flex;flex-direction:column;gap:10px}
.fld{display:flex;flex-direction:column;gap:5px}
.fld>span{
  font:600 9.5px/1 var(--mono);letter-spacing:.19em;text-transform:uppercase;
  color:var(--faint)}
.fld input,.fld textarea,.fld select{width:100%}
.fld textarea{min-height:78px;resize:vertical;line-height:1.5}
.fld.check{flex-direction:row;align-items:center;gap:9px}
.fld.check>span{
  font:14px/1 var(--ui);letter-spacing:0;text-transform:none;color:var(--dim)}
.qrstage{position:sticky;top:70px;display:flex;flex-direction:column;gap:10px}
.sheet{
  width:100%;aspect-ratio:1;border-radius:var(--r-lg);
  border:1px solid var(--edge);display:flex;align-items:center;
  justify-content:center;overflow:hidden;
  background-color:#fff;background-size:16px 16px;background-position:0 0,8px 8px;
  background-image:
    linear-gradient(45deg,#E4E0D8 25%,transparent 25%,transparent 75%,#E4E0D8 75%),
    linear-gradient(45deg,#E4E0D8 25%,transparent 25%,transparent 75%,#E4E0D8 75%)}
.sheet img,.sheet canvas{width:100%;height:100%;object-fit:contain;display:block}
.sheet .empty{
  color:#8B8378;font:11.5px/1.5 var(--mono);padding:20px;text-align:center;
  letter-spacing:.05em}
.qrstat{font:11.5px/1.55 var(--mono);color:var(--faint);text-align:center}
.qrwarn{font-size:12.5px;line-height:1.5;color:var(--amber)}
.logodrop{
  border:1px dashed var(--edge-hi);border-radius:var(--r-lg);padding:14px;
  text-align:center;color:var(--dim);cursor:pointer;font-size:13.5px;
  line-height:1.5;transition:.13s}
.logodrop:hover,.logodrop.hot{
  border-color:var(--steel);color:var(--paper);background:var(--riser)}
.logochip{display:flex;align-items:center;gap:11px}
.logochip img{
  width:44px;height:44px;object-fit:contain;border-radius:var(--r);
  background:#fff;border:1px solid var(--edge);flex:none}
.logochip .nm{flex:1;min-width:0;overflow:hidden;text-overflow:ellipsis;
  white-space:nowrap;font-size:13.5px}

/* ---------- image crop ---------------------------------------------------- */
.cropwrap{
  position:relative;margin:0 auto 12px;line-height:0;width:max-content;
  max-width:100%;border-radius:var(--r-lg);overflow:hidden;
  background:var(--riser);touch-action:none;user-select:none}
.cropwrap img{max-width:100%;display:block;-webkit-user-drag:none}
.cropbox{
  position:absolute;box-sizing:border-box;cursor:move;
  border:1px solid var(--amber);box-shadow:0 0 0 9999px rgba(8,6,5,.66)}
.cropbox .h{
  position:absolute;width:12px;height:12px;background:var(--amber);
  border-radius:1px}
.cropbox .nw{left:-6px;top:-6px;cursor:nwse-resize}
.cropbox .ne{right:-6px;top:-6px;cursor:nesw-resize}
.cropbox .sw{left:-6px;bottom:-6px;cursor:nesw-resize}
.cropbox .se{right:-6px;bottom:-6px;cursor:nwse-resize}

/* ---------- action bar ----------------------------------------------------- */
.actions{
  position:fixed;left:var(--rail-w);right:0;bottom:0;z-index:30;
  display:flex;align-items:center;gap:16px;
  padding:13px 22px;background:var(--veil);backdrop-filter:blur(9px);
  border-top:1px solid var(--edge)}
.status{
  flex:1;min-width:0;font:12px/1.5 var(--mono);color:var(--dim);
  letter-spacing:.02em}
.status.ok{color:var(--ok)}
.status.err{color:var(--rust)}

/* progress — the machine working, so AMBER */
.bar{
  position:fixed;left:0;right:0;top:0;height:2px;background:transparent;
  z-index:60;opacity:0;transition:opacity .2s}
.bar.show{opacity:1}
.bar i{
  display:block;height:100%;width:35%;background:var(--amber);
  box-shadow:0 0 12px #FFB00099;
  animation:slide 1.15s cubic-bezier(.5,0,.5,1) infinite}
@keyframes slide{
  0%{transform:translateX(-100%)}100%{transform:translateX(385%)}}

.warn,.warn2{
  padding:11px 14px;margin-bottom:12px;border-radius:var(--r);
  font-size:13.5px;line-height:1.55;
  background:var(--warn-bg);border:1px solid var(--warn-edge);color:var(--warn-ink)}
.warn2{background:var(--err-bg);border-color:var(--err-edge);color:var(--err-ink)}
.warn:empty,.warn2:empty{display:none}

/* ---------- toasts --------------------------------------------------------- */
.toasts{position:fixed;right:16px;bottom:74px;z-index:80;
  display:flex;flex-direction:column;gap:7px;align-items:flex-end}
.toast{
  padding:9px 14px;border-radius:var(--r);font-size:13.5px;
  background:var(--riser);border:1px solid var(--edge-hi);color:var(--paper);
  box-shadow:0 8px 24px #00000066;
  animation:rise .18s ease}
.toast.ok{border-color:var(--ok-edge);color:var(--ok)}
.toast.err{border-color:var(--err-edge);color:var(--rust)}
@keyframes rise{from{opacity:0;transform:translateY(7px)}}

/* ---------- responsive ------------------------------------------------------ */
@media (max-width:860px){
  .qrgrid{grid-template-columns:1fr}
  .qrstage{position:static}
}
@media (max-width:720px){
  .shell{grid-template-columns:1fr}
  .rail{
    position:static;height:auto;border-right:0;
    border-bottom:1px solid var(--edge);
    display:flex;gap:14px;overflow-x:auto;padding:10px 12px;
    background:var(--bench)}
  .groups{display:none}
  .groups + .tabs{margin-bottom:0}
  .tabs{flex-direction:row;gap:5px}
  .tab{padding:7px 12px;white-space:nowrap;font-size:13.5px}
  .tab::before{display:none}
  .tab.active{background:var(--steel-dim);color:#CFE6F5}
  .stage{padding:16px 14px 120px}
  .actions{left:0;padding:11px 14px}
  .status{font-size:11px}
  .readout{gap:0 20px;padding:11px 14px}
  .ro-cell b{font-size:16px}
  .ytcard{flex-direction:column}
  .ytcard img{width:100%}
}
/* ---------- light chassis --------------------------------------------------
   The same instrument under bench lighting rather than darkroom lighting.
   The readout stays dark in both, because a real display does.
   -------------------------------------------------------------------------- */
:root[data-theme="light"]{
  --void:#E9E4DA; --bench:#F5F1EA; --riser:#FFFFFF;
  --edge:#D3CABB; --edge-hi:#B2A593;
  --paper:#1D1915; --dim:#544C43; --faint:#6E6357;
  --amber:#7D5100; --steel:#255C82; --steel-dim:#D2E2EE;
  --rust:#A2361F; --ok:#376E27;
  --steel-ink:#10394F; --on-steel:#FFFFFF; --on-amber:#FFFFFF;
  --ok-edge:#A9C39C; --err-edge:#D8AEA3; --err-bg:#F7E7E2; --err-ink:#8A2E19;
  --warn-bg:#FBF0D8; --warn-edge:#DCC38A; --warn-ink:#6B4B08;
  --veil:#F5F1EAF2;
}
:root[data-theme="light"] .rail{
  background:linear-gradient(180deg,var(--bench),var(--void) 220px)}
:root[data-theme="light"] .toast{box-shadow:0 8px 24px #00000024}

/* the theme control — a contrast mark, not a sun or a moon */
.theme{
  display:inline-flex;align-items:center;justify-content:center;
  width:28px;height:28px;padding:0;flex:none;cursor:pointer;
  border-radius:var(--r);border:1px solid var(--edge);
  background:var(--void);color:var(--dim);transition:.13s}
.theme:hover{color:var(--paper);border-color:var(--steel);background:var(--riser)}
.theme svg{width:15px;height:15px;display:block;transition:transform .22s ease}
:root[data-theme="light"] .theme svg{transform:rotate(180deg)}

@media (prefers-reduced-motion:reduce){
  *{animation-duration:.001ms!important;transition-duration:.001ms!important}
}
</style></head>
<body>
<header class="masthead">
  <div class="brand">
    <span class="mark" aria-hidden="true">▚</span>
    <h1>Multi Toolkit</h1>
    <button class="theme" id="themeBtn" type="button"
            aria-label="Switch to light theme">
      <svg viewBox="0 0 16 16" aria-hidden="true">
        <circle cx="8" cy="8" r="6.3" fill="none" stroke="currentColor"
                stroke-width="1.5"/>
        <path d="M8 1.7a6.3 6.3 0 0 1 0 12.6z" fill="currentColor"/>
      </svg>
    </button>
    <span class="eyebrow">local workbench</span>
  </div>
  <div class="pills">
    <span class="pill on">127.0.0.1 only</span>
    <span class="pill off" id="pillLO">LibreOffice …</span>
    <span class="pill off" id="pillGS">Ghostscript …</span>
    <span class="pill off" id="pillFF">ffmpeg …</span>
  </div>
</header>

<div class="shell">
  <nav class="rail" aria-label="Tools">
    <div class="groups">
      <div class="group active" data-g="pdf"><b>Documents</b><span>pdf</span></div>
    </div>
    <div class="tabs" data-group="pdf">
      <div class="tab active" data-tab="merge">Merge</div>
      <div class="tab" data-tab="split">Split</div>
      <div class="tab" data-tab="compress">Compress</div>
      <div class="tab" data-tab="convert">Convert</div>
      <div class="tab" data-tab="pdf2md">PDF → MD</div>
    </div>
    <div class="groups">
      <div class="group" data-g="media"><b>Media</b><span>a/v</span></div>
    </div>
    <div class="tabs" data-group="media">
      <div class="tab" data-tab="youtube">YouTube</div>
      <div class="tab" data-tab="reels">Reels</div>
      <div class="tab" data-tab="image">Image</div>
      <div class="tab" data-tab="qr">QR Code</div>
    </div>
  </nav>

  <main class="stage">
  <div class="readout" id="readout" data-state="idle"
       aria-live="polite" aria-label="Measured state"></div>

  <div class="warn" id="warn"></div>

  <div class="drop big" id="drop">
    <span id="dropMain"><b>Drop files here</b> or <b>click to choose</b></span>
    <span class="hint" id="dropHint">Select several at once — Shift-click, Ctrl/Cmd-click, or drag a whole batch in.</span>
    <input type="file" id="file" multiple hidden>
  </div>

  <ul class="files" id="list"></ul>

  <div class="toolbar" id="listTools" style="display:none">
    <button class="btn" id="add">＋ Add more</button>
    <span class="sp"></span>
    <button class="btn danger" id="clear">Clear all</button>
  </div>

  <!-- SPLIT preview -->
  <div id="splitUI" class="hide">
    <div class="chips" id="splitChips">
      <span class="chip active" data-m="cuts">✂ Cut points</span>
      <span class="chip" data-m="extract">Extract pages</span>
      <span class="chip" data-m="reorder">Reorder</span>
      <span class="chip" data-m="each">Every page</span>
      <span class="chip" data-m="every">Every N</span>
      <span class="chip" data-m="ranges">Ranges</span>
    </div>
    <div class="opts" id="splitOpts">
      <label id="everyWrap" class="hide">N <input type="number" id="everyN" value="2" min="1" style="width:64px"></label>
      <label id="rangesWrap" class="hide">Ranges <input type="text" id="ranges" placeholder="1-3,4,5-8" size="16"></label>
      <button class="btn hide" id="resetOrder" style="padding:6px 12px;font-size:13px">Reset order</button>
      <span id="splitTip" style="font-size:12.5px;color:var(--mut)"></span>
    </div>
    <div class="preview" id="preview">
      <div class="pv-head">
        <span class="pv-title" id="pvTitle">Pages</span>
        <span class="pv-sum" id="pvSum"></span>
      </div>
      <div class="pv-note" id="pvNote"></div>
      <div class="grid" id="grid"></div>
    </div>
  </div>

  <!-- option panels -->
  <div class="opts" id="opt-merge">
    <label><input type="checkbox" id="bm" checked> Add a bookmark for each source file</label>
  </div>

  <div id="compressUI" class="hide">
    <div class="presets">
      <div class="pcard" data-p="high">
        <div class="t">High compression</div>
        <div class="d">Smallest file. Images downscaled hard — fine for email &amp; archiving.</div>
        <div class="sz" id="sz-high">—</div>
      </div>
      <div class="pcard active" data-p="balanced">
        <div class="t">Balanced</div>
        <div class="d">Good size cut with little visible loss. The sane default.</div>
        <div class="sz" id="sz-balanced">—</div>
      </div>
      <div class="pcard" data-p="quality">
        <div class="t">High quality</div>
        <div class="d">Light touch on images — for print or figure-heavy papers.</div>
        <div class="sz" id="sz-quality">—</div>
      </div>
      <div class="pcard" data-p="lossless">
        <div class="t">Lossless</div>
        <div class="d">Structure cleanup only. Pixels untouched, guaranteed.</div>
        <div class="sz" id="sz-lossless">—</div>
      </div>
    </div>
    <div id="compNote" style="font-size:12.5px;color:var(--amber);margin:-4px 2px 12px;display:none"></div>
  </div>

  <div class="opts hide" id="opt-convert">
    <label>Conversion
      <select id="route">
        <option value="office2pdf">Office → PDF (Word / PowerPoint / Excel)</option>
        <option value="pdf2docx">PDF → Word (.docx)</option>
        <option value="pdf2pptx">PDF → PowerPoint (.pptx)</option>
        <option value="pdf2img">PDF → Images (PNG / JPG)</option>
        <option value="img2pdf">Images → PDF</option>
      </select>
    </label>
    <label id="imgFmtWrap" class="hide">Format
      <select id="imgFormat"><option value="png">PNG</option><option value="jpg">JPG</option></select>
    </label>
    <label id="dpiWrap" class="hide">DPI <input type="number" id="dpi" value="150" min="72" step="10" style="width:74px"></label>
  </div>

  <!-- YOUTUBE -->
  <div id="ytUI" class="hide">
    <div class="ytbox">
      <div class="ytrow">
        <input type="text" id="ytUrl" placeholder="https://www.youtube.com/watch?v=…"
               autocomplete="off" spellcheck="false">
        <button class="btn" id="ytFetch">Fetch info</button>
      </div>
      <div id="ytCard" class="ytcard hide">
        <img id="ytThumb" alt="">
        <div class="ytmeta">
          <div class="yttitle" id="ytTitle"></div>
          <div class="ytsub" id="ytSub"></div>
          <div class="ytsel">
            <label>Resolution <select id="ytRes"></select></label>
            <label id="ytFmtWrap">Format <select id="ytFmt">
              <option value="qt">QuickTime mp4 — max res, converted (HEVC/H.264)</option>
              <option value="best">Original codec — max quality (VP9/AV1, use VLC/IINA)</option>
              <option value="h264">Native H.264 mp4 — plays anywhere, ≤1080p</option>
            </select></label>
          </div>
        </div>
      </div>
      <div class="opts" id="ytExtra" style="margin:0">
        <label id="ytVertWrap">Vertical 9:16
          <select id="ytVert">
            <option value="off">Keep original framing</option>
            <option value="blur">Blurred backdrop (nothing cropped)</option>
            <option value="crop">Crop to fill</option>
            <option value="pad">Black bars</option>
          </select>
        </label>
        <label id="ytSizeWrap" class="hide">Canvas
          <select id="ytVSize">
            <option value="1080">1080 × 1920 · Reels / Shorts / TikTok</option>
            <option value="720">720 × 1280 · lighter file</option>
            <option value="1350">1080 × 1350 · Instagram feed 4:5</option>
          </select>
        </label>
        <label>Sign-in cookies
          <select id="ytCookies">
            <option value="">None (public posts only)</option>
            <option value="chrome">Chrome</option>
            <option value="firefox">Firefox</option>
            <option value="safari">Safari</option>
            <option value="edge">Edge</option>
            <option value="brave">Brave</option>
            <option value="chromium">Chromium</option>
          </select>
        </label>
      </div>
      <div class="ytnote" id="ytNote">Only download videos you own or have permission to save.</div>
    </div>
  </div>

  <!-- PDF -> MARKDOWN -->
  <div id="mdUI" class="hide">
    <div class="card">
      <h4>For pasting into an LLM</h4>
      <p style="margin:0 0 12px;font-size:14px;color:var(--dim);line-height:1.55">
        Text stays text, so it costs a fraction of the tokens a PDF would.
        Figures are cropped from the page and dropped in at the point they
        appear, each next to its own caption. Pages with no text layer are
        exported whole and flagged, so nothing is silently lost.</p>
      <div class="opts" style="margin:0">
        <label><input type="checkbox" id="mdImages" checked> Extract figures</label>
        <label><input type="checkbox" id="mdTables" checked> Tables as Markdown</label>
        <label><input type="checkbox" id="mdMath" checked> Equations as LaTeX</label>
        <label><input type="checkbox" id="mdHeader" checked> Explain-itself header</label>
        <label><input type="checkbox" id="mdIndex"> Also write _INDEX.md</label>
        <label>Figure quality
          <select id="mdDpi">
            <option value="120">120 dpi — small files</option>
            <option value="170" selected>170 dpi — balanced</option>
            <option value="240">240 dpi — dense figures</option>
            <option value="320">320 dpi — maximum detail</option>
          </select>
        </label>
        <label>Pages
          <input type="text" id="mdPages" placeholder="all" style="width:82px">
        </label>
      </div>
      <div class="ytnote" id="mdNote">Multiple PDFs convert one after another,
        each into its own set of files.</div>
    </div>
  </div>

  <!-- IMAGE -->
  <div id="imgUI" class="hide">
    <div class="ytbox">
      <div class="ytrow">
        <input type="text" id="imgUrl" placeholder="https://…/photo.jpg  (direct link to the image)"
               autocomplete="off" spellcheck="false">
        <button class="btn" id="imgFetch">Load</button>
      </div>
      <div class="logodrop" id="imgDrop">…or <b>drop an image here</b> / click to choose
        <input type="file" id="imgFile" accept="image/*" hidden></div>
    </div>

    <div id="imgEditor" class="hide">
      <div class="card">
        <h4>Crop</h4>
        <div class="cropwrap" id="cropWrap">
          <img id="imgPrev" alt="">
          <div class="cropbox" id="cropBox">
            <span class="h nw"></span><span class="h ne"></span>
            <span class="h sw"></span><span class="h se"></span>
          </div>
        </div>
        <div class="chips" id="arChips" style="margin-bottom:0">
          <span class="chip active" data-ar="free">Free</span>
          <span class="chip" data-ar="1">1:1</span>
          <span class="chip" data-ar="0.8">4:5</span>
          <span class="chip" data-ar="0.5625">9:16</span>
          <span class="chip" data-ar="1.7777778">16:9</span>
          <span class="chip" data-ar="1.5">3:2</span>
          <span class="chip" data-ar="orig">Original</span>
          <span class="chip" data-ar="full">Whole image</span>
        </div>
      </div>

      <div class="opts">
        <label>Width <input type="number" id="outW" min="1" step="1" style="width:92px"></label>
        <label>Height <input type="number" id="outH" min="1" step="1" style="width:92px"></label>
        <label>Format
          <select id="imgFmt">
            <option value="png">PNG — lossless</option>
            <option value="jpg">JPG — small</option>
            <option value="webp">WebP — smallest</option>
          </select>
        </label>
        <label id="imgQWrap" class="hide rng">Quality
          <input type="range" id="imgQ" min="40" max="100" value="90"><b id="imgQVal">90</b>
        </label>
        <label><input type="checkbox" id="imgGray"> Grayscale</label>
        <button class="btn" id="imgRot" style="padding:6px 12px;font-size:13px">Rotate 90°</button>
        <button class="btn" id="imgFlip" style="padding:6px 12px;font-size:13px">Flip</button>
        <button class="btn" id="cropReset" style="padding:6px 12px;font-size:13px">Reset</button>
      </div>
      <div class="ytnote" id="imgInfo"></div>
    </div>
  </div>

  <!-- QR CODE -->
  <div id="qrUI" class="hide">
    <div class="qrgrid">
      <div>
        <div class="card">
          <h4>What should it hold?</h4>
          <div class="chips" id="qrKinds">
            <span class="chip active" data-k="url">Link</span>
            <span class="chip" data-k="text">Text</span>
            <span class="chip" data-k="wifi">Wi-Fi</span>
            <span class="chip" data-k="email">Email</span>
            <span class="chip" data-k="sms">SMS</span>
            <span class="chip" data-k="phone">Phone</span>
            <span class="chip" data-k="vcard">Contact</span>
            <span class="chip" data-k="geo">Location</span>
          </div>
          <div class="qrfields" id="qrFields"></div>
        </div>

        <div class="card">
          <h4>Centre mark</h4>
          <div class="chips" id="qrLogoMode">
            <span class="chip active" data-m="plain">Plain code</span>
            <span class="chip" data-m="logo">With a logo</span>
          </div>
          <div class="logodrop hide" id="qrLogoDrop">
            <b>Drop a logo</b> or click to choose — PNG with a transparent
            background works best
            <input type="file" id="qrLogoFile" accept="image/*" hidden>
          </div>
          <div class="logochip hide" id="qrLogoChip" style="margin-top:11px">
            <img id="qrLogoThumb" alt="">
            <span class="nm" id="qrLogoName"></span>
            <button class="btn danger" id="qrLogoClear"
                    style="padding:5px 11px;font-size:13px">Remove</button>
          </div>
          <div class="opts hide" id="qrLogoOpts" style="margin:11px 0 0">
            <label class="rng">Size
              <input type="range" id="qrLogoPct" min="8" max="34" value="22">
              <b id="qrLogoPctVal">22%</b>
            </label>
            <label>Look
              <select id="qrLogoStyle">
                <option value="original">Original colours</option>
                <option value="silhouette">Silhouette — matches the code</option>
              </select>
            </label>
            <label><input type="checkbox" id="qrPad"> Backing plate</label>
            <label id="qrPadShapeWrap" class="hide">Shape
              <select id="qrPadShape">
                <option value="rounded">Rounded</option>
                <option value="circle">Circle</option>
                <option value="square">Square</option>
              </select>
            </label>
          </div>
        </div>

        <div class="card">
          <h4>Style</h4>
          <div class="opts" style="margin:0;padding:0;border:0;background:none">
            <label>Modules
              <select id="qrStyle">
                <option value="square">Square</option>
                <option value="dots">Dots</option>
                <option value="rounded">Rounded</option>
              </select>
            </label>
            <label>Foreground <input type="color" id="qrFg" value="#000000" class="swatch"
                                    style="width:44px"></label>
            <label>Background <input type="color" id="qrBg" value="#ffffff" class="swatch"
                                    style="width:44px"></label>
            <label><input type="checkbox" id="qrTransparent"> Transparent</label>
            <label>Error correction
              <select id="qrEc">
                <option value="L">L — 7%</option>
                <option value="M">M — 15%</option>
                <option value="Q">Q — 25%</option>
                <option value="H" selected>H — 30% (needed for logos)</option>
              </select>
            </label>
            <label>Size
              <select id="qrSize">
                <option value="512">512 px</option>
                <option value="1024" selected>1024 px</option>
                <option value="2048">2048 px — print</option>
                <option value="4096">4096 px — poster</option>
              </select>
            </label>
            <label>Quiet zone
              <input type="number" id="qrBorder" value="4" min="0" max="16" style="width:64px">
            </label>
            <label>File
              <select id="qrFmt"><option value="png">PNG</option><option value="jpg">JPG</option></select>
            </label>
          </div>
        </div>
      </div>

      <div class="qrstage">
        <div class="sheet" id="qrSheet">
          <span class="empty" id="qrEmpty">Your code appears here as you type</span>
          <img id="qrImg" alt="QR preview" style="display:none">
        </div>
        <div class="qrstat" id="qrStat"></div>
        <div class="qrwarn" id="qrWarn"></div>
      </div>
    </div>
  </div>

  <div class="actions">
    <div class="status" id="status">Ready — drop a file to begin.</div>
    <button class="primary" id="go">Merge &amp; Download</button>
  </div>

  </main>
</div>

<div class="bar" id="bar"><i></i></div>
<div class="toasts" id="toasts"></div>

<script>
const $=id=>document.getElementById(id);
const GC=['#8b7bff','#46d6a4','#ffc66b','#5ec8ff','#ff8ad4'];
let files=[];                 // {id,name,size,b64,kind,token,pages,cover,thumbs}
let tab='merge';
let splitId=null;             // file id previewed in Split
let splitMode='cuts';
let cuts=new Set();           // cut AFTER page p
let sel=new Set();            // selected pages (extract)
let order=null;               // page order for Reorder mode
let compPreset='balanced';
let compPrev=null;            // {sig, data} from /compress_preview
let ytInfo=null, ytJob=null, ytPoll=null, ytBusy=false;
let qrLogo=null, qrKind='url', qrResult=null, qrSeq=0, qrTimer=null;
let imgState={token:null,name:'',w:0,h:0,crop:null,ar:null,rotate:0,flip:false,
              url:'',b64:null};
const PDF_TABS=['merge','split','compress','convert','pdf2md'];
const VIDEO_TABS=['youtube','reels'];
const GROUP_OF={merge:'pdf',split:'pdf',compress:'pdf',convert:'pdf',
  pdf2md:'pdf',youtube:'media',reels:'media',image:'media',qr:'media'};
const GO_LABEL={merge:'Merge & Download',split:'Split & Download',
  compress:'Compress & Download',convert:'Convert & Download',
  pdf2md:'Convert to Markdown',
  youtube:'Download Video',reels:'Download Reel',image:'Export Image',
  qr:'Download QR Code'};
const lastTab={pdf:'merge',media:'youtube'};
let caps={libreoffice:false, ghostscript:false, ffmpeg:false};
let nextId=1;

const ACCEPT={
  merge:'.pdf,application/pdf', split:'.pdf,application/pdf', compress:'.pdf,application/pdf',
  office2pdf:'.doc,.docx,.ppt,.pptx,.xls,.xlsx,.odt,.odp,.ods,.rtf',
  pdf2docx:'.pdf', pdf2pptx:'.pdf', pdf2img:'.pdf',
  img2pdf:'.png,.jpg,.jpeg,.bmp,.gif,.tif,.tiff,.webp,image/*'
};
const IMG_EXT={png:'image/png',jpg:'image/jpeg',jpeg:'image/jpeg',gif:'image/gif',
  bmp:'image/bmp',webp:'image/webp',tif:'image/tiff',tiff:'image/tiff'};

function human(n){const u=['B','KB','MB','GB'];let i=0;while(n>=1024&&i<3){n/=1024;i++;}
  return (i===0?n:n.toFixed(1))+' '+u[i];}
function status(m,c){const s=$('status');s.textContent=m;s.className='status'+(c?' '+c:'');}
function toast(m,c){const t=document.createElement('div');t.className='toast'+(c?' '+c:'');
  t.textContent=m;$('toasts').appendChild(t);setTimeout(()=>t.remove(),3200);}
function kindOf(name){const e=name.split('.').pop().toLowerCase();
  if(e==='pdf')return 'pdf'; if(IMG_EXT[e])return 'img'; return 'other';}
function currentRoute(){ return tab==='convert' ? $('route').value : tab; }
function fref(f,full){ return (!full && f.token) ? {token:f.token,name:f.name}
                                                 : {name:f.name,data:f.b64}; }

/* ---------- adding files ---------- */
function bufToB64(buf){let bin='';const b=new Uint8Array(buf),c=0x8000;
  for(let i=0;i<b.length;i+=c)bin+=String.fromCharCode.apply(null,b.subarray(i,i+c));
  return btoa(bin);}

async function addFiles(fl){
  let added=0;
  for(const file of [...fl]){
    if(files.some(f=>f.name===file.name && f.size===file.size)) continue;
    const buf=await file.arrayBuffer();
    const f={id:nextId++,name:file.name,size:file.size,b64:bufToB64(buf),
             kind:kindOf(file.name),token:null,pages:null,cover:null,thumbs:null};
    files.push(f); added++;
    if(f.kind==='pdf') inspectCover(f);          // async: token + pages + cover
  }
  if(added && tab==='split' && splitId===null){
    const p=files.find(f=>f.kind==='pdf'); if(p){splitId=p.id; resetSplitState();}
  }
  render();
  if(added){toast(`Added ${added} file${added>1?'s':''}`,'ok');
    if(tab==='compress')refreshCompressPreview();}
  else toast('Duplicates skipped');
}

async function inspectCover(f){
  try{
    const res=await fetch('/inspect',{method:'POST',headers:{'Content-Type':'application/json'},
      body:JSON.stringify({file:{name:f.name,data:f.b64},thumbs:'cover'})});
    if(!res.ok) return;
    const j=await res.json();
    f.token=j.token; f.pages=j.pages; f.cover=j.cover;
    render();
    if(tab==='split' && splitId===f.id) refreshPreview();
  }catch(e){/* offline-ish failure: ops will still work with raw data */}
}

async function ensureThumbs(f){
  if(f.thumbs) return f;
  let res=await fetch('/inspect',{method:'POST',headers:{'Content-Type':'application/json'},
    body:JSON.stringify({file:fref(f,false),thumbs:'all'})});
  if(res.status===409){
    f.token=null;
    res=await fetch('/inspect',{method:'POST',headers:{'Content-Type':'application/json'},
      body:JSON.stringify({file:fref(f,true),thumbs:'all'})});
  }
  if(!res.ok) throw new Error(await res.text());
  const j=await res.json();
  f.token=j.token; f.pages=j.pages; f.thumbs=j.thumbs||[]; f.cover=j.cover||f.cover;
  return f;
}

/* ---------- file list ---------- */
function coverSrc(f){
  if(f.kind==='pdf'&&f.cover) return 'data:image/jpeg;base64,'+f.cover;
  if(f.kind==='img'){const e=f.name.split('.').pop().toLowerCase();
    return 'data:'+(IMG_EXT[e]||'image/png')+';base64,'+f.b64;}
  return null;
}

function render(){
  const list=$('list'); list.innerHTML='';
  const pickMode = tab==='split';
  files.forEach((f,i)=>{
    const li=document.createElement('li');
    li.className='row'+(pickMode&&f.kind==='pdf'?' pickable':'')
                 +(pickMode&&f.id===splitId?' picked':'');
    li.draggable=true; li.dataset.id=f.id;
    const src=coverSrc(f);
    const thumb=src?`<img class="thumb" src="${src}" alt="">`
      :`<span class="thumb">${(f.name.split('.').pop()||'?').toUpperCase().slice(0,4)}</span>`;
    const pages=f.kind==='pdf' ? (f.pages!=null?` · ${f.pages} pg`:' · …') : '';
    li.innerHTML=`<span class="grip" title="Drag to reorder">⋮⋮</span>${thumb}
      <span class="meta"><span class="nm" title="${f.name}">${f.name}</span>
        <span class="sub">${human(f.size)}${pages}</span></span>
      <span class="pickdot">PREVIEWING</span>
      <button data-up="${i}" ${i===0?'disabled':''} title="Move up">↑</button>
      <button data-down="${i}" ${i===files.length-1?'disabled':''} title="Move down">↓</button>
      <button class="del" data-del="${i}" title="Remove">✕</button>`;
    list.appendChild(li);
  });
  const has=!!files.length, usesList=PDF_TABS.includes(tab);
  $('listTools').style.display=(has&&usesList)?'flex':'none';
  const drop=$('drop');
  drop.classList.toggle('big',!has); drop.classList.toggle('slim',has);
  drop.classList.toggle('hide',!usesList); list.classList.toggle('hide',!usesList);
  if(VIDEO_TABS.includes(tab))   $('go').disabled = !ytInfo||ytBusy;
  else if(tab==='qr')            $('go').disabled = !qrResult;
  else if(tab==='image')         $('go').disabled = !imgState.token;
  else if(tab==='pdf2md')        $('go').disabled = !files.some(f=>f.kind==='pdf');
  else                           $('go').disabled = !has;
  if(tab==='split') refreshSplitVisibility();
  updateReadout();
}

$('list').addEventListener('click',e=>{
  const t=e.target;
  if(t.dataset.del!==undefined){
    const f=files[+t.dataset.del];
    files.splice(+t.dataset.del,1);
    if(f.id===splitId){const p=files.find(x=>x.kind==='pdf');
      splitId=p?p.id:null; resetSplitState();}
    render(); if(tab==='split') refreshPreview();
    if(tab==='compress') refreshCompressPreview();
  }
  else if(t.dataset.up!==undefined){const i=+t.dataset.up;
    [files[i-1],files[i]]=[files[i],files[i-1]];render();}
  else if(t.dataset.down!==undefined){const i=+t.dataset.down;
    [files[i+1],files[i]]=[files[i],files[i+1]];render();}
  else if(tab==='split'){
    const li=e.target.closest('.row'); if(!li) return;
    const f=files.find(x=>x.id===+li.dataset.id);
    if(f&&f.kind==='pdf'&&f.id!==splitId){splitId=f.id;resetSplitState();render();refreshPreview();}
  }
});

/* drag to reorder */
let dragEl=null;
$('list').addEventListener('dragstart',e=>{
  dragEl=e.target.closest('.row'); if(!dragEl)return;
  dragEl.classList.add('dragging'); e.dataTransfer.effectAllowed='move';
});
$('list').addEventListener('dragover',e=>{
  e.preventDefault();
  const over=e.target.closest('.row');
  if(!over||!dragEl||over===dragEl) return;
  const r=over.getBoundingClientRect();
  const before=(e.clientY-r.top)<r.height/2;
  over.parentNode.insertBefore(dragEl, before?over:over.nextSibling);
});
$('list').addEventListener('drop',e=>e.preventDefault());
$('list').addEventListener('dragend',()=>{
  if(!dragEl) return;
  dragEl.classList.remove('dragging'); dragEl=null;
  const order=[...$('list').children].map(li=>+li.dataset.id);
  files.sort((a,b)=>order.indexOf(a.id)-order.indexOf(b.id));
  render();
});

/* reorder-mode drag within the page grid */
let pgDrag=null;
$('grid').addEventListener('dragstart',e=>{
  if(splitMode!=='reorder')return;
  pgDrag=e.target.closest('.pg'); if(!pgDrag)return;
  pgDrag.classList.add('dragging'); e.dataTransfer.effectAllowed='move';
});
$('grid').addEventListener('dragover',e=>{
  if(!pgDrag)return; e.preventDefault();
  const over=e.target.closest('.pg');
  if(!over||over===pgDrag)return;
  const r=over.getBoundingClientRect();
  const before=(e.clientX-r.left)<r.width/2;
  over.parentNode.insertBefore(pgDrag, before?over:over.nextSibling);
});
$('grid').addEventListener('drop',e=>{if(pgDrag)e.preventDefault();});
$('grid').addEventListener('dragend',()=>{
  if(!pgDrag)return;
  pgDrag.classList.remove('dragging'); pgDrag=null;
  order=[...$('grid').querySelectorAll('.pg')].map(el=>+el.dataset.p);
  drawGrid();
});
$('resetOrder').onclick=()=>{order=null;drawGrid();};

$('add').onclick=()=>$('file').click();
$('drop').onclick=()=>$('file').click();
$('file').onchange=e=>{addFiles([...e.target.files]);e.target.value='';};
$('clear').onclick=()=>{if(files.length&&confirm('Remove all files?')){
  files=[];splitId=null;resetSplitState();compPrev=null;render();refreshPreview();
  refreshCompressPreview();status('Cleared.');}};

['dragenter','dragover'].forEach(ev=>document.body.addEventListener(ev,e=>{
  e.preventDefault();
  if([...(e.dataTransfer?.types||[])].includes('Files'))$('drop').classList.add('hot');}));
['dragleave','drop'].forEach(ev=>document.body.addEventListener(ev,e=>{
  e.preventDefault();if(ev==='drop'||e.target===document.body)$('drop').classList.remove('hot');}));
document.body.addEventListener('drop',e=>{
  if(!PDF_TABS.includes(tab)) return;   // other tabs have their own drop zones
  if(e.dataTransfer.files.length)addFiles(e.dataTransfer.files);});

/* ---------- tabs ---------- */
function updateAccept(){
  // The element is #file; this tab only takes PDFs.
  $('file').accept = (tab==='pdf2md') ? '.pdf,application/pdf'
                                      : (ACCEPT[currentRoute()]||'');
}
function selectTab(name){
  tab=name; const g=GROUP_OF[name]; lastTab[g]=name;
  document.querySelectorAll('.tab').forEach(t=>
    t.classList.toggle('active',t.dataset.tab===name));
  document.querySelectorAll('.group').forEach(x=>
    x.classList.toggle('active',x.dataset.g===g));
  ['merge','convert'].forEach(t=>$('opt-'+t)&&$('opt-'+t).classList.toggle('hide',t!==name));
  $('compressUI').classList.toggle('hide',name!=='compress');
  $('mdUI').classList.toggle('hide',name!=='pdf2md');
  $('ytUI').classList.toggle('hide',!VIDEO_TABS.includes(name));
  $('imgUI').classList.toggle('hide',name!=='image');
  $('qrUI').classList.toggle('hide',name!=='qr');
  if(name!=='split') $('splitUI').classList.add('hide');
  $('go').textContent=GO_LABEL[name];
  if(VIDEO_TABS.includes(name)) applyVideoMode();
  if(name==='compress') refreshCompressPreview();
  if(name==='split'&&splitId===null){const p=files.find(f=>f.kind==='pdf');
    if(p){splitId=p.id;resetSplitState();}}
  updateAccept(); refreshWarn(); render();
  if(name==='split') refreshPreview();
  if(name==='qr') qrRefresh();
}
document.querySelectorAll('.tab').forEach(el=>el.onclick=()=>selectTab(el.dataset.tab));
document.querySelectorAll('.group').forEach(el=>el.onclick=()=>selectTab(lastTab[el.dataset.g]));

/* ---------- split preview ---------- */
function resetSplitState(){cuts=new Set();sel=new Set();order=null;}
function refreshSplitVisibility(){
  $('splitUI').classList.toggle('hide', tab!=='split' || !splitId);
}

document.querySelectorAll('#splitChips .chip').forEach(c=>c.onclick=()=>{
  document.querySelectorAll('#splitChips .chip').forEach(x=>x.classList.remove('active'));
  c.classList.add('active'); splitMode=c.dataset.m;
  $('everyWrap').classList.toggle('hide',splitMode!=='every');
  $('rangesWrap').classList.toggle('hide',splitMode!=='ranges');
  $('resetOrder').classList.toggle('hide',splitMode!=='reorder');
  drawGrid();
});
$('everyN').oninput=drawGrid;
$('ranges').oninput=drawGrid;

function parseClientRanges(spec,total){
  const out=[];
  for(let chunk of String(spec).split(',')){
    chunk=chunk.trim(); if(!chunk)continue;
    let a,b;
    if(chunk.includes('-')){const[x,y]=chunk.split('-');
      a=x.trim()?parseInt(x):1; b=y.trim()?parseInt(y):total;}
    else a=b=parseInt(chunk);
    if(isNaN(a)||isNaN(b))continue;
    a=Math.max(1,a);b=Math.min(total,b);
    if(a<=b)out.push([a,b]);
  }
  return out;
}

function groupsFor(total){
  // returns array: group index per page (0-based), or -1 = not included
  const g=new Array(total).fill(-1);
  if(splitMode==='cuts'){
    let gi=0;
    for(let p=1;p<=total;p++){g[p-1]=gi; if(cuts.has(p))gi++;}
  }else if(splitMode==='each'){
    for(let p=0;p<total;p++)g[p]=p;
  }else if(splitMode==='every'){
    const n=Math.max(1,parseInt($('everyN').value)||1);
    for(let p=0;p<total;p++)g[p]=Math.floor(p/n);
  }else if(splitMode==='ranges'){
    parseClientRanges($('ranges').value,total).forEach(([a,b],i)=>{
      for(let p=a;p<=b;p++)g[p-1]=i;});
  }else if(splitMode==='extract'){
    for(const p of sel)g[p-1]=0;
  }
  return g;
}

async function refreshPreview(){
  refreshSplitVisibility();
  if(tab!=='split'||!splitId)return;
  const f=files.find(x=>x.id===splitId);
  if(!f){refreshSplitVisibility();return;}
  $('pvTitle').textContent=f.name;
  if(!f.thumbs){
    $('grid').innerHTML='<div style="color:var(--mut);font-size:13px;padding:8px">Rendering page previews…</div>';
    try{await ensureThumbs(f);}catch(err){
      $('grid').innerHTML='<div style="color:var(--bad);font-size:13px;padding:8px">Could not render previews: '
        +err.message+'</div>';return;}
  }
  drawGrid();
}

function drawGrid(){
  const f=files.find(x=>x.id===splitId);
  if(tab!=='split'||!f||!f.thumbs)return;
  const total=f.pages||f.thumbs.length;
  const grid=$('grid');
  const notes={
    cuts:'Click between pages to place a ✂ cut. Each tinted group becomes its own PDF.',
    extract:'Click pages to select them. Selected pages are pulled into one PDF (in order).',
    reorder:'Drag pages into a new order; hover a page and hit ✕ to remove it. Output is one rebuilt PDF.',
    each:'Every page becomes its own single-page PDF.',
    every:'Pages are chunked into groups of N.',
    ranges:'Type ranges like 1-3,4,5-8. Dimmed pages are left out.'
  };
  $('pvNote').textContent=notes[splitMode];

  if(splitMode==='reorder'){
    if(!order)order=Array.from({length:total},(_,i)=>i+1);
    grid.innerHTML='';
    for(const p of order){
      const card=document.createElement('div');
      card.className='pg reorderable'; card.draggable=true; card.dataset.p=p;
      const img=f.thumbs[p-1];
      card.innerHTML=(img?`<img src="data:image/jpeg;base64,${img}" alt="page ${p}">`
                         :`<div class="ph">${p}</div>`)
        +`<div class="lbl">p.${p}</div><span class="rm" title="Remove page">✕</span>`;
      card.querySelector('.rm').onclick=e=>{e.stopPropagation();
        order=order.filter(x=>x!==p);drawGrid();};
      grid.appendChild(card);
    }
    const moved=order.some((p,i)=>p!==i+1)||order.length!==total;
    $('pvSum').textContent=`${order.length} of ${total} pages → 1 PDF`
      +(moved?'':' (original order)');
    return;
  }

  const g=groupsFor(total);
  grid.innerHTML='';
  for(let p=1;p<=total;p++){
    const card=document.createElement('div');
    const gi=g[p-1];
    card.className='pg'+(splitMode==='extract'?' selectable':'')
      +((gi<0)?' dim':'')
      +((splitMode==='extract'&&sel.has(p))?' sel':'');
    if(gi>=0&&splitMode!=='extract')card.style.borderColor=GC[gi%GC.length];
    const img=f.thumbs[p-1];
    card.innerHTML=(img?`<img src="data:image/jpeg;base64,${img}" alt="page ${p}">`
                       :`<div class="ph">${p}</div>`)
      +`<div class="lbl">p.${p}${gi>=0&&splitMode!=='extract'&&splitMode!=='each'?' · G'+(gi+1):''}</div>`;
    if(splitMode==='extract')card.onclick=()=>{sel.has(p)?sel.delete(p):sel.add(p);drawGrid();};
    grid.appendChild(card);
    if(p<total){
      const cut=document.createElement('div');
      const isBoundary=(g[p-1]!==g[p]&&g[p-1]>=0&&g[p]>=0);
      cut.className='cut'+(splitMode==='cuts'?' clickable':'')
        +(splitMode==='cuts'&&cuts.has(p)?' on':'')
        +(splitMode!=='cuts'&&isBoundary?' boundary':'');
      cut.innerHTML='<div class="ln"></div><span class="sc">✂</span>';
      if(splitMode==='cuts'){
        cut.title='Cut after page '+p;
        cut.onclick=()=>{cuts.has(p)?cuts.delete(p):cuts.add(p);drawGrid();};
      }
      grid.appendChild(cut);
    }
  }
  // summary
  let sum='';
  if(splitMode==='cuts')sum=`${cuts.size} cut${cuts.size===1?'':'s'} → ${cuts.size+1} file${cuts.size?'s':''}`;
  else if(splitMode==='extract')sum=`${sel.size} page${sel.size===1?'':'s'} selected → 1 PDF`;
  else if(splitMode==='each')sum=`${total} pages → ${total} files`;
  else if(splitMode==='every'){const n=Math.max(1,parseInt($('everyN').value)||1);
    sum=`${total} pages → ${Math.ceil(total/n)} files`;}
  else if(splitMode==='ranges'){const r=parseClientRanges($('ranges').value,total);
    sum=`${r.length} range${r.length===1?'':'s'} → ${r.length} file${r.length===1?'':'s'}`;}
  $('pvSum').textContent=sum;
}

function cutsToRanges(total){
  const pts=[...cuts].sort((a,b)=>a-b);
  const out=[];let start=1;
  for(const c of pts){out.push(start+'-'+c);start=c+1;}
  out.push(start+'-'+total);
  return out.join(',');
}

/* ---------- compress presets ---------- */
document.querySelectorAll('.pcard').forEach(c=>c.onclick=()=>{
  document.querySelectorAll('.pcard').forEach(x=>x.classList.remove('active'));
  c.classList.add('active'); compPreset=c.dataset.p;
});

function compSig(){return files.filter(f=>f.kind==='pdf').map(f=>f.id).join(',');}

async function refreshCompressPreview(){
  if(tab!=='compress')return;
  const pdfs=files.filter(f=>f.kind==='pdf');
  const sig=compSig();
  if(!pdfs.length){compPrev=null;
    document.querySelectorAll('.pcard .sz').forEach(s=>{s.textContent='\u2014';s.className='sz';});
    return;}
  if(compPrev&&compPrev.sig===sig&&compPrev.data){paintSizes();return;}
  if(compPrev&&compPrev.sig===sig&&compPrev.pending)return;   // already estimating
  compPrev={sig,pending:true,data:null};
  document.querySelectorAll('.pcard .sz').forEach(s=>{s.textContent='estimating\u2026';s.className='sz wait';});
  try{
    let res=await fetch('/compress_preview',{method:'POST',
      headers:{'Content-Type':'application/json'},
      body:JSON.stringify({files:pdfs.map(f=>fref(f,false))})});
    if(res.status===409){
      pdfs.forEach(f=>f.token=null);
      res=await fetch('/compress_preview',{method:'POST',
        headers:{'Content-Type':'application/json'},
        body:JSON.stringify({files:pdfs.map(f=>fref(f,true))})});
    }
    if(!res.ok)throw new Error(await res.text());
    const data=await res.json();
    if(compSig()!==sig)return;          // files changed while estimating
    compPrev={sig,data};
    paintSizes();
  }catch(err){
    compPrev=null;
    document.querySelectorAll('.pcard .sz').forEach(s=>{
      s.textContent='preview failed \u2014 download still works';s.className='sz wait';});
  }
}

function paintSizes(){
  if(!compPrev||!compPrev.data)return;
  const orig=compPrev.data.orig;
  for(const[k,v]of Object.entries(compPrev.data.presets)){
    const el=$('sz-'+k); if(!el)continue;
    const pct=Math.round(100*v.total/Math.max(1,orig));
    el.textContent=`${human(orig)} \u2192 ${human(v.total)} (${pct}%)`;
    el.className='sz';
  }
  // If every preset lands in the same place, say why instead of looking broken.
  const pp=compPrev.data.presets;
  const note=$('compNote');
  const spread=Math.abs(pp.lossless.total-pp.high.total)/Math.max(1,orig);
  if(spread<0.05){
    note.style.display='block';
    note.textContent='Presets all land within a few percent \u2014 this PDF is mostly text/vector with '
      +'little raster imagery to recompress.'
      +(caps.ghostscript?'':' Installing Ghostscript unlocks much stronger whole-file compression '
        +'(fonts, streams, downsampling): macOS \u201cbrew install ghostscript\u201d, Windows '
        +'\u201cwinget install ArtifexSoftware.GhostScript\u201d \u2014 then restart this app.');
  }else{
    note.style.display='none';
  }
}

/* ---------- option toggles ---------- */
$('route').onchange=()=>{const r=$('route').value;
  $('imgFmtWrap').classList.toggle('hide',r!=='pdf2img');
  $('dpiWrap').classList.toggle('hide',!(r==='pdf2img'||r==='pdf2pptx'));
  updateAccept(); refreshWarn();};

function refreshWarn(){
  const w=$('warn');
  if(currentRoute()==='office2pdf' && !caps.libreoffice){
    w.style.display='block';
    w.textContent='⚠ LibreOffice not detected — needed for Office → PDF. Install from libreoffice.org/download, then restart this app.';
  } else { w.style.display='none'; }
}

/* ---------- run ---------- */
function payload(full){
  const route=currentRoute();
  const refs=files.map(f=>fref(f,full));
  if(tab==='merge') return {ep:'/merge', body:{files:refs,bookmarks:$('bm').checked}};
  if(tab==='split'){
    const f=files.find(x=>x.id===splitId)||files.find(x=>x.kind==='pdf');
    if(!f) throw new Error('Add a PDF to split.');
    const total=f.pages||1;
    const body={file:fref(f,full)};
    if(splitMode==='cuts'){
      if(!cuts.size) throw new Error('Place at least one ✂ cut between pages (or pick another mode).');
      body.mode='ranges'; body.ranges=cutsToRanges(total);
    }else if(splitMode==='extract'){
      if(!sel.size) throw new Error('Click some pages to select them first.');
      body.mode='extract'; body.pages=[...sel];
    }else if(splitMode==='reorder'){
      if(!order||!order.length) throw new Error('No pages left \u2014 hit Reset order or restore pages.');
      body.mode='reorder'; body.pages=order;
    }else if(splitMode==='every'){body.mode='every'; body.every_n=+$('everyN').value;}
    else if(splitMode==='ranges'){body.mode='ranges'; body.ranges=$('ranges').value;}
    else body.mode='each';
    return {ep:'/split', body};
  }
  if(tab==='pdf2md'){
    const pdfs=files.filter(f=>f.kind==='pdf');
    let pf=null, pt=null;
    const m=/^\s*(\d+)\s*[-:]\s*(\d+)\s*$/.exec($('mdPages').value.trim());
    if(m){ pf=+m[1]; pt=+m[2]; }
    else if(/^\s*\d+\s*$/.test($('mdPages').value)){ pf=pt=+$('mdPages').value; }
    return {ep:'/pdf2md', body:{file:pdfs.length
        ? refs[files.indexOf(pdfs[0])] : refs[0],
      dpi:+$('mdDpi').value, images:$('mdImages').checked,
      tables:$('mdTables').checked, math:$('mdMath').checked,
      header:$('mdHeader').checked, index:$('mdIndex').checked,
      page_from:pf, page_to:pt}};
  }
  if(tab==='compress'){
    // If the preview already compressed everything, just download the cached bytes.
    if(!full && compPrev && compPrev.data && compPrev.sig===compSig()){
      const pf=compPrev.data.presets[compPreset];
      return {ep:'/download', body:{files:pf.files.map(x=>({token:x.token,name:x.name})),
        info:'Compressed ('+compPreset+'): '+human(compPrev.data.orig)+' \u2192 '+human(pf.total)}};
    }
    return {ep:'/compress', body:{files:refs,preset:compPreset}};
  }
  if(tab==='image'){
    if(!imgState.token) throw new Error('Load an image first.');
    const c=imgState.crop||{x:0,y:0,w:imgState.w,h:imgState.h};
    return {ep:'/img_process', body:{
      file:{token:imgState.token,name:imgState.name},
      crop:{x:Math.round(c.x),y:Math.round(c.y),
            w:Math.round(c.w),h:Math.round(c.h)},
      out_w:+$('outW').value||null, out_h:+$('outH').value||null,
      fmt:$('imgFmt').value, quality:+$('imgQ').value,
      rotate:imgState.rotate, flip:imgState.flip,
      grayscale:$('imgGray').checked}};
  }
  return {ep:'/convert', body:{files:refs,route,
      img_format:$('imgFormat').value,dpi:+$('dpi').value}};
}

async function send(full){
  const {ep,body}=payload(full);
  return fetch(ep,{method:'POST',headers:{'Content-Type':'application/json'},
    body:JSON.stringify(body)});
}

/* ---------- youtube ---------- */
function fmtDur(s){s=Math.round(s||0);const h=Math.floor(s/3600),m=Math.floor(s%3600/60),x=s%60;
  return (h?h+':'+String(m).padStart(2,'0'):m)+':'+String(x).padStart(2,'0');}
function fmtEta(s){if(s==null)return'';s=Math.round(s);
  return s>=60?Math.floor(s/60)+'m '+(s%60)+'s':s+'s';}

async function ytFetchInfo(){
  const url=$('ytUrl').value.trim();
  if(!url){toast('Paste a video URL first');return;}
  ytInfo=null;$('ytCard').classList.add('hide');render();
  $('ytFetch').disabled=true;$('bar').classList.add('show');status('Fetching video info…');
  try{
    const res=await fetch('/yt_info',{method:'POST',
      headers:{'Content-Type':'application/json'},body:JSON.stringify({url})});
    if(!res.ok)throw new Error(await res.text());
    ytInfo=await res.json();
    $('ytTitle').textContent=ytInfo.title;
    $('ytSub').textContent=[ytInfo.uploader,ytInfo.duration?fmtDur(ytInfo.duration):'']
      .filter(Boolean).join(' · ');
    if(ytInfo.thumbnail){$('ytThumb').src=ytInfo.thumbnail;$('ytThumb').style.display='';}
    else $('ytThumb').style.display='none';
    const sel=$('ytRes');sel.innerHTML='';
    const hs=ytInfo.heights&&ytInfo.heights.length?ytInfo.heights:[null];
    hs.forEach((h,i)=>{const o=document.createElement('option');
      o.value=h||'';o.textContent=h?h+'p'+(i===0?'  (max)':''):'Best available';
      sel.appendChild(o);});
    const note=$('ytNote');
    $('ytFmtWrap').classList.toggle('hide',!ytInfo.has_ffmpeg);
    if(!ytInfo.has_ffmpeg){note.className='ytnote warn2';
      note.textContent='⚠ ffmpeg not found — YouTube caps single-stream downloads around 720p. '
        +'Install ffmpeg for full 1080p/4K (macOS “brew install ffmpeg”, Windows '
        +'“winget install Gyan.FFmpeg”, Linux “apt install ffmpeg”), then restart this app.';
    }else{note.className='ytnote';
      note.textContent='YouTube stores >1080p only as VP9/AV1, which QuickTime can\u2019t play \u2014 '
        +'pick \u201cQuickTime mp4\u201d to auto-convert, or \u201cOriginal codec\u201d if you use VLC/IINA. '
        +'Only download videos you own or have permission to save.';}
    $('ytCard').classList.remove('hide');
    status('Ready — pick a resolution and hit Download Video.','ok');
  }catch(err){status('Failed: '+err.message,'err');toast('Could not read that URL','err');}
  finally{$('ytFetch').disabled=false;$('bar').classList.remove('show');render();}
}
$('ytFetch').onclick=ytFetchInfo;
$('ytUrl').addEventListener('keydown',e=>{if(e.key==='Enter')ytFetchInfo();});
$('ytUrl').addEventListener('input',()=>{   // new URL invalidates old info
  if(ytInfo){ytInfo=null;$('ytCard').classList.add('hide');render();}
});

async function ytGo(){
  if(!ytInfo||ytBusy)return;
  ytBusy=true;render();$('bar').classList.add('show');status('Starting download…');
  try{
    const res=await fetch('/yt_start',{method:'POST',
      headers:{'Content-Type':'application/json'},
      body:JSON.stringify({url:$('ytUrl').value.trim(),height:+$('ytRes').value||null,
        compat:$('ytFmtWrap').classList.contains('hide')?'best':$('ytFmt').value,
        vertical:vertMode(), vsize:$('ytVSize').value,
        cookies:$('ytCookies').value})});
    if(!res.ok)throw new Error(await res.text());
    ytJob=(await res.json()).id;
    ytPoll=setInterval(ytCheck,600);
  }catch(err){ytDone();status('Failed: '+err.message,'err');toast('Download failed','err');}
}
function ytDone(){clearInterval(ytPoll);ytPoll=null;ytJob=null;ytBusy=false;
  $('bar').classList.remove('show');render();}
async function ytCheck(){
  if(!ytJob)return;
  try{
    const res=await fetch('/yt_progress?id='+ytJob);
    if(!res.ok)throw new Error(await res.text());
    const j=await res.json();
    if(j.status==='error'){ytDone();
      status('Failed: '+(j.error||'unknown error'),'err');toast('Download failed','err');return;}
    if(j.status==='done'){
      const a=document.createElement('a');a.href='/yt_file?id='+ytJob;
      a.download=j.name||'video.mp4';document.body.appendChild(a);a.click();a.remove();
      ytDone();
      const det=j.detail?('\nActual quality: '+j.detail):'';
      status('Downloaded: '+(j.name||'video')+det,'ok');
      toast('Downloaded'+(j.detail?' \u00b7 '+j.detail:''),'ok');return;}
    let line=(j.msg||'Working')+'…';
    if(j.pct!=null)line+='  '+j.pct+'%';
    if(j.speed)line+='  ·  '+human(j.speed)+'/s';
    if(j.eta!=null)line+='  ·  ETA '+fmtEta(j.eta);
    if(j.rate)line+='  ·  '+j.rate;
    status(line);
  }catch(e){/* transient poll failure — keep trying */}
}

$('go').onclick=async()=>{
  if(VIDEO_TABS.includes(tab)){ytGo();return;}
  if(tab==='qr'){qrDownload();return;}
  if(tab!=='image' && !files.length) return;
  $('go').disabled=true; $('bar').classList.add('show'); status('Working…');
  try{
    let res=await send(false);
    if(res.status===409){          // cache expired (server restarted) — resend bytes
  if(tab==='image'){ await imgReload(); res=await send(false); }
      else{
        files.forEach(f=>f.token=null); compPrev=null;
        files.filter(f=>f.kind==='pdf').forEach(inspectCover);
        res=await send(true);
      }
    }
    if(!res.ok){throw new Error(await res.text()||('HTTP '+res.status));}
    const fname=decodeURIComponent(res.headers.get('X-Filename')||'output');
    const info=decodeURIComponent(res.headers.get('X-Info')||'');
    const blob=await res.blob();
    const url=URL.createObjectURL(blob);
    const a=document.createElement('a');a.href=url;a.download=fname;
    document.body.appendChild(a);a.click();a.remove();URL.revokeObjectURL(url);
    status((info||('Saved '+fname))+'\nDownloaded: '+fname,'ok');
    toast('Downloaded '+fname,'ok');
  }catch(err){status('Failed: '+err.message,'err');toast('Failed — see status','err');}
  finally{$('bar').classList.remove('show');render();}
};


/* ---------- video mode (YouTube vs Reels) ---------- */
function applyVideoMode(){
  const reels = tab==='reels';
  $('ytUrl').placeholder = reels
    ? 'https://www.instagram.com/reel/…  (also TikTok, Shorts, X)'
    : 'https://www.youtube.com/watch?v=…';
  // 9:16 reframing is a Reels concern — it only clutters the YouTube tab.
  // The select keeps its value so a Reels choice survives a trip to YouTube.
  $('ytVertWrap').classList.toggle('hide', !reels);
  if(reels && !$('ytVert').dataset.touched) $('ytVert').value = 'blur';
  syncVertSize();
}
function vertMode(){ return tab==='reels' ? $('ytVert').value : 'off'; }
function syncVertSize(){
  $('ytSizeWrap').classList.toggle('hide', vertMode()==='off');
}
$('ytVert').onchange=()=>{ $('ytVert').dataset.touched='1'; syncVertSize(); };

/* ---------- QR code ---------- */
const QR_KINDS={
  url:  [['text','Link','example.com','text']],
  text: [['text','Text','Anything you want the code to carry','area']],
  wifi: [['ssid','Network name (SSID)','MyNetwork','text'],
         ['password','Password','','text'],
         ['security','Security','','sel:WPA/WPA2|WEP|nopass/No password'],
         ['hidden','Hidden network','','check']],
  email:[['to','To','name@example.com','text'],
         ['subject','Subject','','text'],
         ['body','Message','','area']],
  sms:  [['phone','Phone number','+1 555 0100','text'],
         ['message','Message','','area']],
  phone:[['phone','Phone number','+1 555 0100','text']],
  vcard:[['name','Full name','Ada Lovelace','text'],
         ['org','Organisation','','text'],['title','Job title','','text'],
         ['phone','Phone','','text'],['email','Email','','text'],
         ['url','Website','','text']],
  geo:  [['lat','Latitude','35.7796','text'],['lon','Longitude','-78.6382','text']]
};

function qrBuildFields(){
  const box=$('qrFields'); box.innerHTML='';
  for(const [key,label,ph,type] of QR_KINDS[qrKind]){
    const w=document.createElement('label');
    w.className='fld'+(type==='check'?' check':'');
    const id='qrf_'+key;
    if(type==='check'){
      w.innerHTML=`<input type="checkbox" id="${id}"><span>${label}</span>`;
    }else if(type==='area'){
      w.innerHTML=`<span>${label}</span><textarea id="${id}" placeholder="${ph}"></textarea>`;
    }else if(type.startsWith('sel:')){
      const opts=type.slice(4).split('|').map(o=>{
        const [v,t]=o.split('/'); return `<option value="${v}">${t||v}</option>`;}).join('');
      w.innerHTML=`<span>${label}</span><select id="${id}">${opts}</select>`;
    }else{
      w.innerHTML=`<span>${label}</span><input type="text" id="${id}" placeholder="${ph}">`;
    }
    box.appendChild(w);
    const el=$(id);
    el.addEventListener(type==='check'||type.startsWith('sel:')?'change':'input',qrSoon);
  }
}
function qrFieldValues(){
  const out={};
  for(const [key,,,type] of QR_KINDS[qrKind]){
    const el=$('qrf_'+key); if(!el) continue;
    out[key]= type==='check' ? (el.checked?'1':'') : el.value;
  }
  return out;
}
document.querySelectorAll('#qrKinds .chip').forEach(c=>c.onclick=()=>{
  document.querySelectorAll('#qrKinds .chip').forEach(x=>x.classList.remove('active'));
  c.classList.add('active'); qrKind=c.dataset.k; qrBuildFields(); qrRefresh();
});

function qrSoon(){ clearTimeout(qrTimer); qrTimer=setTimeout(qrRefresh,260); }
['qrStyle','qrFg','qrBg','qrTransparent','qrEc','qrSize','qrBorder','qrFmt',
 'qrLogoStyle','qrPadShape'].forEach(id=>$(id).addEventListener('change',qrRefresh));
$('qrBorder').addEventListener('input',qrSoon);
$('qrLogoPct').addEventListener('input',()=>{
  $('qrLogoPctVal').textContent=$('qrLogoPct').value+'%'; qrSoon();});
$('qrPad').addEventListener('change',()=>{
  $('qrPadShapeWrap').classList.toggle('hide',!$('qrPad').checked); qrRefresh();});
$('qrTransparent').addEventListener('change',()=>{
  $('qrBg').disabled=$('qrTransparent').checked;});

/* plain vs logo — an explicit choice, so a plain code needs no upload */
function qrSetLogoMode(mode){
  const wantLogo = mode==='logo';
  document.querySelectorAll('#qrLogoMode .chip').forEach(c=>
    c.classList.toggle('active', c.dataset.m===mode));
  $('qrLogoDrop').classList.toggle('hide', !wantLogo || !!qrLogo);
  $('qrLogoChip').classList.toggle('hide', !wantLogo || !qrLogo);
  $('qrLogoOpts').classList.toggle('hide', !wantLogo || !qrLogo);
  if(!wantLogo && qrLogo){ qrLogo=null; $('qrLogoThumb').removeAttribute('src'); }
  qrRefresh();
}
document.querySelectorAll('#qrLogoMode .chip').forEach(c=>
  c.onclick=()=>qrSetLogoMode(c.dataset.m));

/* logo picking */
function qrSetLogo(file){
  if(!file) return;
  if(!/^image\//.test(file.type)){toast('That is not an image','err');return;}
  const r=new FileReader();
  r.onload=()=>{
    qrLogo={name:file.name,data:String(r.result).split(',')[1],url:String(r.result)};
    $('qrLogoThumb').src=qrLogo.url; $('qrLogoName').textContent=file.name;
    if($('qrEc').value!=='H'){ $('qrEc').value='H';
      toast('Error correction raised to H for the logo'); }
    qrSetLogoMode('logo');
  };
  r.readAsDataURL(file);
}
$('qrLogoDrop').onclick=()=>$('qrLogoFile').click();
$('qrLogoFile').onchange=e=>{qrSetLogo(e.target.files[0]);e.target.value='';};
['dragenter','dragover'].forEach(ev=>$('qrLogoDrop').addEventListener(ev,e=>{
  e.preventDefault();e.stopPropagation();$('qrLogoDrop').classList.add('hot');}));
['dragleave','drop'].forEach(ev=>$('qrLogoDrop').addEventListener(ev,e=>{
  e.preventDefault();e.stopPropagation();$('qrLogoDrop').classList.remove('hot');}));
$('qrLogoDrop').addEventListener('drop',e=>qrSetLogo(e.dataTransfer.files[0]));
$('qrLogoClear').onclick=()=>{
  qrLogo=null; $('qrLogoThumb').removeAttribute('src');
  qrSetLogoMode('logo');   // stay in logo mode, ready for another file
};

async function qrRefresh(){
  if(tab!=='qr') return;
  const fields=qrFieldValues();
  const empty=!Object.entries(fields).some(([k,v])=>v&&v.trim());
  if(empty){
    qrResult=null; $('qrImg').style.display='none'; $('qrEmpty').style.display='';
    $('qrStat').textContent=''; $('qrWarn').textContent=''; render(); return;
  }
  const seq=++qrSeq;
  try{
    const res=await fetch('/qr',{method:'POST',
      headers:{'Content-Type':'application/json'},
      body:JSON.stringify({
        kind:qrKind, fields,
        ec:$('qrEc').value, target_px:+$('qrSize').value,
        border:Math.max(0,Math.min(16,+$('qrBorder').value||0)),
        fg:$('qrFg').value,
        bg:$('qrTransparent').checked?'transparent':$('qrBg').value,
        style:$('qrStyle').value,
        logo_data:qrLogo?qrLogo.data:null,
        logo_pct:+$('qrLogoPct').value,
        logo_style:$('qrLogoStyle').value,
        pad:$('qrPad').checked, pad_shape:$('qrPadShape').value, pad_pct:8,
        fmt:$('qrFmt').value})});
    if(seq!==qrSeq) return;                    // a newer keystroke won
    if(!res.ok){throw new Error(await res.text());}
    const j=await res.json();
    if(seq!==qrSeq) return;
    qrResult=j;
    $('qrImg').src='data:'+j.mime+';base64,'+j.image;
    $('qrImg').style.display=''; $('qrEmpty').style.display='none';
    $('qrStat').textContent=`v${j.version} · ${j.modules}×${j.modules} modules · `
      +`EC ${j.ec} · ${j.size}px · ${human(j.bytes)} · ${j.chars} chars`;
    $('qrWarn').textContent=j.warn||'';
    status('QR code ready.','ok');
  }catch(err){
    if(seq!==qrSeq) return;
    qrResult=null; $('qrImg').style.display='none'; $('qrEmpty').style.display='';
    $('qrStat').textContent=''; $('qrWarn').textContent='';
    status('QR failed: '+err.message,'err');
  }
  render();
}

function qrDownload(){
  if(!qrResult){toast('Nothing to download yet');return;}
  const bin=atob(qrResult.image);
  const arr=new Uint8Array(bin.length);
  for(let i=0;i<bin.length;i++)arr[i]=bin.charCodeAt(i);
  const blob=new Blob([arr],{type:qrResult.mime});
  const url=URL.createObjectURL(blob);
  const stem=(qrFieldValues().text||qrKind).replace(/^https?:\/\//,'')
    .replace(/[^\w.\-]+/g,'_').slice(0,40)||'qr';
  const a=document.createElement('a');
  a.href=url; a.download=stem+'_qr.'+($('qrFmt').value==='jpg'?'jpg':'png');
  document.body.appendChild(a);a.click();a.remove();URL.revokeObjectURL(url);
  toast('Downloaded','ok');
  status('Saved '+a.download+' — scan it once before you print it.','ok');
}

/* ---------- image: load, crop, export ---------- */
function imgAR(){
  const v=imgState.ar;
  if(v===null||v==='free')return null;
  if(v==='orig')return imgState.w/imgState.h;
  if(v==='full')return null;
  return parseFloat(v);
}
function cropScale(){ const im=$('imgPrev'); return (im.clientWidth||1)/(imgState.w||1); }
function drawCrop(){
  const c=imgState.crop; if(!c)return;
  const s=cropScale(), b=$('cropBox');
  b.style.left=(c.x*s)+'px';  b.style.top=(c.y*s)+'px';
  b.style.width=(c.w*s)+'px'; b.style.height=(c.h*s)+'px';
  syncOutFields();
}
function syncOutFields(){
  setTimeout(updateReadout,0);
  const c=imgState.crop; if(!c)return;
  let w=Math.round(c.w), h=Math.round(c.h);
  if(imgState.rotate===90||imgState.rotate===270){const t=w;w=h;h=t;}
  $('outW').value=w; $('outH').value=h;
  $('imgInfo').textContent=
    `Source ${imgState.w}×${imgState.h} · crop ${Math.round(c.w)}×${Math.round(c.h)}`
    +` at (${Math.round(c.x)}, ${Math.round(c.y)})`
    +(imgState.rotate?` · rotated ${imgState.rotate}°`:'')
    +(imgState.flip?' · flipped':'');
}
function fitCrop(ar){
  const W=imgState.w,H=imgState.h;
  if(!ar){ imgState.crop={x:0,y:0,w:W,h:H}; return; }
  let w=W, h=w/ar;
  if(h>H){ h=H; w=h*ar; }
  imgState.crop={x:(W-w)/2, y:(H-h)/2, w, h};
}
document.querySelectorAll('#arChips .chip').forEach(c=>c.onclick=()=>{
  document.querySelectorAll('#arChips .chip').forEach(x=>x.classList.remove('active'));
  c.classList.add('active'); imgState.ar=c.dataset.ar;
  if(c.dataset.ar==='full'){ imgState.crop={x:0,y:0,w:imgState.w,h:imgState.h}; }
  else if(c.dataset.ar!=='free'){ fitCrop(imgAR()); }
  drawCrop();
});
$('cropReset').onclick=()=>{
  imgState.rotate=0; imgState.flip=false;
  imgState.crop={x:0,y:0,w:imgState.w,h:imgState.h};
  document.querySelectorAll('#arChips .chip').forEach((x,i)=>
    x.classList.toggle('active',i===0));
  imgState.ar='free'; drawCrop();
};
$('imgRot').onclick=()=>{imgState.rotate=(imgState.rotate+90)%360;syncOutFields();
  toast('Rotation applies on export');};
$('imgFlip').onclick=()=>{imgState.flip=!imgState.flip;syncOutFields();
  toast('Flip applies on export');};
$('imgFmt').onchange=()=>$('imgQWrap').classList.toggle('hide',$('imgFmt').value==='png');
$('imgQ').oninput=()=>$('imgQVal').textContent=$('imgQ').value;
$('outW').oninput=()=>{
  const c=imgState.crop; if(!c)return;
  const w=+$('outW').value; if(w>0) $('outH').value=Math.max(1,Math.round(w*c.h/c.w));
};
$('outH').oninput=()=>{
  const c=imgState.crop; if(!c)return;
  const h=+$('outH').value; if(h>0) $('outW').value=Math.max(1,Math.round(h*c.w/c.h));
};

let cropDrag=null;
$('cropWrap').addEventListener('pointerdown',e=>{
  if(!imgState.crop)return;
  const cls=[...e.target.classList];
  const hand=cls.includes('h')?['nw','ne','sw','se'].find(h=>cls.includes(h)):null;
  if(!hand && e.target.id!=='cropBox') return;
  e.preventDefault();
  $('cropWrap').setPointerCapture(e.pointerId);
  cropDrag={hand,sx:e.clientX,sy:e.clientY,start:{...imgState.crop}};
});
$('cropWrap').addEventListener('pointermove',e=>{
  if(!cropDrag)return;
  const s=cropScale(), st=cropDrag.start, W=imgState.w, H=imgState.h;
  const dx=(e.clientX-cropDrag.sx)/s, dy=(e.clientY-cropDrag.sy)/s;
  const MIN=16;
  if(!cropDrag.hand){
    imgState.crop={x:Math.max(0,Math.min(W-st.w,st.x+dx)),
                   y:Math.max(0,Math.min(H-st.h,st.y+dy)),w:st.w,h:st.h};
  }else{
    let x0=st.x, y0=st.y, x1=st.x+st.w, y1=st.y+st.h;
    if(cropDrag.hand.includes('w')) x0=st.x+dx; else x1=st.x+st.w+dx;
    if(cropDrag.hand.includes('n')) y0=st.y+dy; else y1=st.y+st.h+dy;
    x0=Math.max(0,Math.min(x0,x1-MIN)); x1=Math.min(W,Math.max(x1,x0+MIN));
    y0=Math.max(0,Math.min(y0,y1-MIN)); y1=Math.min(H,Math.max(y1,y0+MIN));
    let c={x:x0,y:y0,w:x1-x0,h:y1-y0};
    const ar=imgAR();
    if(ar){                                   // keep the locked ratio exact
      let w=c.w, h=w/ar;
      if(h>H){h=H;w=h*ar;}
      if(cropDrag.hand.includes('w')) c.x=x1-w;
      if(cropDrag.hand.includes('n')) c.y=y1-h;
      c.w=w; c.h=h;
      c.x=Math.max(0,Math.min(c.x,W-w)); c.y=Math.max(0,Math.min(c.y,H-h));
    }
    imgState.crop=c;
  }
  drawCrop();
});
['pointerup','pointercancel'].forEach(ev=>
  $('cropWrap').addEventListener(ev,()=>{cropDrag=null;}));
window.addEventListener('resize',()=>{ if(tab==='image'&&imgState.crop) drawCrop(); });

async function imgLoad(body){
  $('bar').classList.add('show'); status('Loading image…');
  try{
    const res=await fetch('/img_fetch',{method:'POST',
      headers:{'Content-Type':'application/json'},body:JSON.stringify(body)});
    if(!res.ok) throw new Error(await res.text());
    const j=await res.json();
    imgState.token=j.token; imgState.name=j.name;
    imgState.w=j.width; imgState.h=j.height;
    imgState.crop={x:0,y:0,w:j.width,h:j.height};
    imgState.ar='free'; imgState.rotate=0; imgState.flip=false;
    document.querySelectorAll('#arChips .chip').forEach((x,i)=>
      x.classList.toggle('active',i===0));
    const im=$('imgPrev');
    im.onload=()=>drawCrop();
    im.src='data:image/jpeg;base64,'+j.preview;
    $('imgEditor').classList.remove('hide');
    drawCrop();
    status(`${j.name} — ${j.width}×${j.height} ${j.format}, ${human(j.bytes)}. `
      +'Drag the box to crop.','ok');
  }catch(err){
    status('Could not load that image: '+err.message,'err');
    toast('Image load failed','err');
  }finally{$('bar').classList.remove('show');render();}
}
async function imgReload(){
  if(imgState.b64) return imgLoad({file:{name:imgState.name,data:imgState.b64}});
  if(imgState.url) return imgLoad({url:imgState.url});
}
$('imgFetch').onclick=()=>{
  const u=$('imgUrl').value.trim();
  if(!u){toast('Paste an image URL first');return;}
  imgState.url=u; imgState.b64=null; imgLoad({url:u});
};
$('imgUrl').addEventListener('keydown',e=>{if(e.key==='Enter')$('imgFetch').click();});
$('imgDrop').onclick=()=>$('imgFile').click();
function imgTakeFile(file){
  if(!file)return;
  if(!/^image\//.test(file.type)){toast('That is not an image','err');return;}
  const r=new FileReader();
  r.onload=()=>{
    const b64=String(r.result).split(',')[1];
    imgState.b64=b64; imgState.url=''; imgState.name=file.name;
    imgLoad({file:{name:file.name,data:b64}});
  };
  r.readAsDataURL(file);
}
$('imgFile').onchange=e=>{imgTakeFile(e.target.files[0]);e.target.value='';};
['dragenter','dragover'].forEach(ev=>$('imgDrop').addEventListener(ev,e=>{
  e.preventDefault();e.stopPropagation();$('imgDrop').classList.add('hot');}));
['dragleave','drop'].forEach(ev=>$('imgDrop').addEventListener(ev,e=>{
  e.preventDefault();e.stopPropagation();$('imgDrop').classList.remove('hot');}));
$('imgDrop').addEventListener('drop',e=>imgTakeFile(e.dataTransfer.files[0]));


/* ---------- theme ---------- */
function applyTheme(t){
  document.documentElement.setAttribute('data-theme',t);
  document.documentElement.style.colorScheme=t;
  $('themeBtn').setAttribute('aria-label',
    'Switch to '+(t==='dark'?'light':'dark')+' theme');
}
applyTheme(document.documentElement.getAttribute('data-theme')||'dark');
$('themeBtn').onclick=()=>{
  const next=document.documentElement.getAttribute('data-theme')==='dark'
    ? 'light' : 'dark';
  applyTheme(next);
  try{ localStorage.setItem('mtk-theme',next); }catch(e){}
};

/* ---------- the readout ---------- */
function setReadout(state, cells){
  const el=$('readout'); if(!el) return;
  el.dataset.state=state;
  el.innerHTML=cells.map(c=>
    '<span class="ro-cell"><b>'+c[1]+'</b><i>'+c[0]+'</i></span>').join('');
}
function clock(sec){
  sec=Math.round(sec||0);
  const h=Math.floor(sec/3600), m=Math.floor(sec%3600/60), s2=sec%60;
  const p=n=>String(n).padStart(2,'0');
  return h ? h+':'+p(m)+':'+p(s2) : m+':'+p(s2);
}
function updateReadout(){
  if(tab==='qr'){
    if(!qrResult) return setReadout('idle',[['qr code','awaiting input']]);
    return setReadout('live',[
      ['version','v'+qrResult.version],
      ['grid',qrResult.modules+'×'+qrResult.modules],
      ['correction','EC-'+qrResult.ec],
      ['output',qrResult.size+' px'],
      ['payload',qrResult.chars+' ch']]);
  }
  if(tab==='pdf2md'){
    const pdfs=files.filter(f=>f.kind==='pdf');
    if(!pdfs.length) return setReadout('idle',[['pdf \u2192 md','add a pdf']]);
    const pg=pdfs.reduce((a,f)=>a+(f.pages||0),0);
    const cells=[['queued',pdfs.length]];
    if(pg) cells.push(['pages',pg]);
    cells.push(['figures',$('mdImages').checked?'cropped':'off'],
               ['header',$('mdHeader').checked?'on':'off']);
    return setReadout('live',cells);
  }
  if(tab==='image'){
    if(!imgState.token) return setReadout('idle',[['image','nothing loaded']]);
    const c=imgState.crop||{w:imgState.w,h:imgState.h};
    return setReadout('live',[
      ['source',imgState.w+'×'+imgState.h],
      ['crop',Math.round(c.w)+'×'+Math.round(c.h)],
      ['export',($('outW').value||'—')+'×'+($('outH').value||'—')],
      ['format',$('imgFmt').value.toUpperCase()]]);
  }
  if(VIDEO_TABS.includes(tab)){
    if(!ytInfo) return setReadout('idle',[['video','paste a link']]);
    const cells=[['length',clock(ytInfo.duration)],
                 ['target',($('ytRes').value||'max')+(($('ytRes').value)?'p':'')]];
    if(tab==='reels'&&vertMode()!=='off')
      cells.push(['reframe',$('ytVSize').value==='1350'?'1080×1350':
        ($('ytVSize').value==='720'?'720×1280':'1080×1920')]);
    return setReadout('live',cells);
  }
  if(!files.length) return setReadout('idle',[['files','none loaded']]);
  const bytes=files.reduce((a,f)=>a+(f.size||0),0);
  const pages=files.reduce((a,f)=>
    a+(typeof f.pages==='number'?f.pages:0),0);
  const cells=[['files',files.length]];
  if(pages) cells.push(['pages',pages]);
  cells.push(['total',human(bytes)]);
  return setReadout('live',cells);
}

/* ---------- boot ---------- */
fetch('/capabilities').then(r=>r.json()).then(c=>{
  caps=c;
  $('pillLO').textContent='LibreOffice '+(c.libreoffice?'✓':'—');
  $('pillLO').className='pill '+(c.libreoffice?'on':'off');
  $('pillGS').textContent='Ghostscript '+(c.ghostscript?'✓':'—');
  $('pillGS').className='pill '+(c.ghostscript?'on':'off');
  $('pillFF').textContent='ffmpeg '+(c.ffmpeg?'✓':'—');
  $('pillFF').className='pill '+(c.ffmpeg?'on':'off');
  refreshWarn();
}).catch(()=>{});
qrBuildFields(); updateAccept(); selectTab('merge');
</script></body></html>
"""


def main():
    port = _free_port()
    url = f"http://127.0.0.1:{port}/"
    server = ThreadingHTTPServer(("127.0.0.1", port), Handler)
    lo = "yes" if office_binary() else "no (Office→PDF disabled until installed)"
    gs = "yes (best compression)" if gs_binary() else "no (pure-Python compression used)"
    ff = ("yes (full-res video + 9:16 reframing)" if ffmpeg_binary()
          else "no (video capped ~720p, no vertical reframing)")
    print("=" * 60)
    print("  Multi Toolkit — running locally")
    print("  PDF: merge · split · compress · convert · pdf→md")
    print("  Media: youtube · reels · image · qr code")
    print(f"  Open:          {url}")
    print(f"  LibreOffice:   {lo}")
    print(f"  Ghostscript:   {gs}")
    print(f"  ffmpeg:        {ff}")
    print("  Ctrl-C here to stop.")
    print("=" * 60)
    threading.Timer(0.7, lambda: webbrowser.open(url)).start()
    try:
        server.serve_forever()
    except KeyboardInterrupt:
        print("\nShutting down.")
        server.shutdown()


if __name__ == "__main__":
    if len(sys.argv) > 1 and sys.argv[1] in ("md", "pdf2md"):
        md_cli(sys.argv[2:])          # headless batch conversion
    else:
        main()                        # normal: serve the UI
